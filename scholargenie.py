#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ScholarGenie — AI Research Intelligence
Interactive REPL terminal. Type commands, see results inline.
"""

import os
import sys
import io
import json
import sqlite3
import re
import time
import webbrowser
from pathlib import Path
from typing import List, Dict, Optional
import warnings
warnings.filterwarnings("ignore")
os.environ.setdefault("TF_CPP_MIN_LOG_LEVEL", "3")
os.environ.setdefault("TF_ENABLE_ONEDNN_OPTS", "0")

# ── Windows encoding fix ───────────────────────────────────────────────────────
if sys.platform == "win32":
    try:
        if hasattr(sys.stdout, "buffer"):
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
        if hasattr(sys.stderr, "buffer"):
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")
    except Exception:
        pass

# ── Auto-install packages ─────────────────────────────────────────────────────
def _ensure():
    import subprocess, importlib.util
    needed = {"arxiv": "arxiv", "requests": "requests", "pptx": "python-pptx", "rich": "rich"}
    missing = [pkg for mod, pkg in needed.items() if importlib.util.find_spec(mod) is None]
    if missing:
        print(f"Installing: {', '.join(missing)} ...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-q"] + missing)
_ensure()

# ── Rich ──────────────────────────────────────────────────────────────────────
from rich.console import Console
from rich.table import Table
from rich.panel import Panel
from rich.text import Text
from rich.rule import Rule
from rich.columns import Columns
from rich.padding import Padding
from rich.progress import Progress, SpinnerColumn, BarColumn, TextColumn, TimeElapsedColumn
from rich import box

console = Console()

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = Path(__file__).parent
DATA   = BASE / "data"
PPTX   = DATA / "presentations"
DB_PATH = DATA / "scholargenie.db"
for d in [DATA, PPTX]:
    d.mkdir(parents=True, exist_ok=True)

# ── Session state ─────────────────────────────────────────────────────────────
SESSION: Dict = {
    "results":  [],   # last search results
    "query":    "",
}

# ══════════════════════════════════════════════════════════════════════════════
# DATABASE
# ══════════════════════════════════════════════════════════════════════════════

class DB:
    def __init__(self):
        self.conn = sqlite3.connect(str(DB_PATH), check_same_thread=False)
        self._init()

    def _init(self):
        self.conn.executescript("""
            CREATE TABLE IF NOT EXISTS papers (
                id TEXT PRIMARY KEY, title TEXT, authors TEXT,
                abstract TEXT, year INTEGER, url TEXT, source TEXT,
                citations INTEGER DEFAULT 0, created_at TEXT DEFAULT (datetime('now'))
            );
            CREATE TABLE IF NOT EXISTS summaries (
                paper_id TEXT PRIMARY KEY, tldr TEXT, short_summary TEXT,
                full_summary TEXT, keypoints TEXT
            );
        """)
        self.conn.commit()

    def upsert(self, p):
        self.conn.execute(
            "INSERT OR REPLACE INTO papers VALUES(?,?,?,?,?,?,?,?,datetime('now'))",
            (p["id"], p["title"], json.dumps(p.get("authors", [])),
             p.get("abstract", ""), p.get("year", 0),
             p.get("url", ""), p.get("source", ""), p.get("citations", 0))
        )
        self.conn.commit()

    def save_summary(self, pid, s):
        self.conn.execute(
            "INSERT OR REPLACE INTO summaries VALUES(?,?,?,?,?)",
            (pid, s.get("tldr",""), s.get("short",""), s.get("full",""), json.dumps(s.get("keypoints",[])))
        )
        self.conn.commit()

    def get_summary(self, pid):
        r = self.conn.execute(
            "SELECT tldr,short_summary,full_summary,keypoints FROM summaries WHERE paper_id=?", (pid,)
        ).fetchone()
        if not r: return None
        return {"tldr": r[0], "short": r[1], "full": r[2], "keypoints": json.loads(r[3] or "[]")}

    def all_papers(self):
        rows = self.conn.execute(
            "SELECT id,title,authors,abstract,year,url,source,citations FROM papers ORDER BY created_at DESC"
        ).fetchall()
        return [{"id":r[0],"title":r[1],"authors":json.loads(r[2] or "[]"),
                 "abstract":r[3],"year":r[4],"url":r[5],"source":r[6],"citations":r[7]} for r in rows]

    def get(self, pid):
        r = self.conn.execute(
            "SELECT id,title,authors,abstract,year,url,source,citations FROM papers WHERE id=?", (pid,)
        ).fetchone()
        if not r: return None
        return {"id":r[0],"title":r[1],"authors":json.loads(r[2] or "[]"),
                "abstract":r[3],"year":r[4],"url":r[5],"source":r[6],"citations":r[7]}

    def count(self):
        return self.conn.execute("SELECT COUNT(*) FROM papers").fetchone()[0]

db = DB()


# ══════════════════════════════════════════════════════════════════════════════
# PAPER FINDER
# ══════════════════════════════════════════════════════════════════════════════

def _arxiv(query, n):
    try:
        import arxiv
        results = []
        for r in arxiv.Search(query=query, max_results=n*2,
                              sort_by=arxiv.SortCriterion.Relevance).results():
            results.append({
                "id": r.entry_id.split("/")[-1],
                "title": r.title.strip(),
                "authors": [a.name for a in r.authors],
                "abstract": r.summary.strip(),
                "year": r.published.year,
                "url": r.pdf_url or "",
                "source": "arXiv",
                "citations": 0,
            })
        return results
    except Exception as e:
        console.print(f"  [dim red]arXiv: {e}[/dim red]")
        return []

def _semantic(query, n):
    try:
        import requests
        r = requests.get(
            "https://api.semanticscholar.org/graph/v1/paper/search",
            params={"query": query, "limit": n*2,
                    "fields": "title,authors,abstract,year,citationCount,externalIds"},
            headers={"x-api-key": os.getenv("SEMANTIC_SCHOLAR_API_KEY", "")},
            timeout=12,
        )
        if r.status_code != 200: return []
        out = []
        for item in r.json().get("data", []):
            pid = item.get("paperId", "unknown")
            ext = item.get("externalIds") or {}
            if ext.get("ArXiv"): pid = ext["ArXiv"]
            out.append({
                "id": pid,
                "title": (item.get("title") or "Unknown").strip(),
                "authors": [a.get("name","") for a in item.get("authors",[])],
                "abstract": (item.get("abstract") or "No abstract.").strip(),
                "year": item.get("year") or 0,
                "citations": item.get("citationCount") or 0,
                "url": f"https://semanticscholar.org/paper/{pid}",
                "source": "Semantic Scholar",
            })
        return out
    except Exception as e:
        console.print(f"  [dim red]Semantic Scholar: {e}[/dim red]")
        return []

def find_papers(query: str, n: int) -> List[Dict]:
    all_p = _arxiv(query, n) + _semantic(query, n)
    seen, unique = set(), []
    for p in all_p:
        key = re.sub(r"\W+", "", p["title"].lower())[:60]
        if key not in seen and len(key) > 4:
            seen.add(key); unique.append(p)
    return unique[:n]


# ══════════════════════════════════════════════════════════════════════════════
# SUMMARIZER
# ══════════════════════════════════════════════════════════════════════════════

def summarize(paper: Dict) -> Dict:
    abstract = paper.get("abstract", "").strip()
    if not abstract or len(abstract) < 30:
        return {"tldr": "No abstract.", "short": "No abstract.", "full": "No abstract.", "keypoints": []}

    sents = [s.strip() for s in re.split(r"(?<=[.!?])\s+", abstract) if s.strip()]
    kws   = ["propose","present","introduce","novel","new","show","demonstrate",
             "achieve","outperform","state-of-the-art","we ","our ","this paper"]

    scored = [(sum(2 for k in kws if k in s.lower()) + min(len(s.split())/5,3), s) for s in sents]
    scored.sort(reverse=True)
    tldr  = " ".join(scored[0][1].split()[:60])
    if len(scored[0][1].split()) > 60: tldr += "..."

    markers = ["we propose","we present","we introduce","we show","our approach",
               "our method","we achieve","results show","outperform","novel","first"]
    kpoints = [s for s in sents if any(m in s.lower() for m in markers)][:5]
    if not kpoints:
        step = max(1, len(sents)//5)
        kpoints = sents[::step][:5]

    full = abstract
    try:
        from transformers import pipeline
        pipe = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6",
                        max_length=180, min_length=50, truncation=True)
        full = pipe(abstract[:1024], do_sample=False)[0]["summary_text"]
    except Exception:
        pass

    return {"tldr": tldr, "short": " ".join(sents[:3]), "full": full, "keypoints": kpoints}


# ══════════════════════════════════════════════════════════════════════════════
# PRESENTATION GENERATOR  — 14-slide academic layout
# ══════════════════════════════════════════════════════════════════════════════

_PPTX_STOP = {
    "the","a","an","of","in","for","and","or","to","is","are","with","that",
    "this","we","our","on","at","by","from","as","be","it","its","also","but",
    "not","have","has","was","were","using","used","based","which","paper",
    "show","can","when","their","they","been","more","than","these","such",
    "both","each","after","before","other","into","through","during","over",
}

def _pptx_content(paper: Dict, summary: Dict) -> Dict:
    """Extract structured content for all 14 slides."""
    abstract = paper.get("abstract", "")
    title    = paper.get("title", "")
    authors  = paper.get("authors", [])

    sents = [s.strip() for s in re.split(r"(?<=[.!?])\s+", abstract) if len(s.strip()) > 20]

    def pick(keywords, n=3):
        hits = [s for s in sents if any(k in s.lower() for k in keywords)]
        return hits[:n] if hits else sents[:min(n, len(sents))]

    text  = (title + " " + abstract).lower()
    words = re.findall(r"\b[a-z][a-z]{4,}\b", text)
    freq: Dict[str, int] = {}
    for w in words:
        if w not in _PPTX_STOP:
            freq[w] = freq.get(w, 0) + 1
    modules = sorted(freq, key=freq.get, reverse=True)[:8]

    auth_str = ", ".join(authors[:3])
    if len(authors) > 3:
        auth_str += " et al."

    return {
        "title":      title,
        "authors":    auth_str,
        "year":       paper.get("year", ""),
        "source":     paper.get("source", ""),
        "url":        paper.get("url", ""),
        "tldr":       summary.get("tldr", ""),
        "keypoints":  summary.get("keypoints", []),
        "abstract":   abstract,
        "intro":      sents[:3],
        "problems":   pick(["however","challenge","limitation","problem","difficult",
                            "lack","existing","current","traditional","previous"], 3),
        "objectives": pick(["aim","objective","goal","propose","present","introduce",
                            "develop","design","investigate","focus"], 3),
        "method":     pick(["propose","present","method","approach","framework",
                            "algorithm","model","architecture","system","technique"], 4),
        "results":    pick(["achieve","result","accuracy","performance","improve",
                            "outperform","show","demonstrate","experiment",
                            "state-of-the-art","evaluate"], 4),
        "future":     pick(["future","extend","plan","further","next","potential",
                            "could","will","limitation","open"], 3),
        "modules":    modules,
    }


def make_pptx(paper: Dict, summary: Dict) -> Optional[str]:
    try:
        from pptx import Presentation as PPT
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        prs = PPT()
        c   = _pptx_content(paper, summary)
        m   = c["modules"]

        NAVY = RGBColor(0x1e, 0x1b, 0x4b)

        def _bold_title(sl):
            try:
                for para in sl.shapes.title.text_frame.paragraphs:
                    para.font.bold = True
                    para.font.color.rgb = NAVY
            except Exception:
                pass

        def add_slide(title_text: str, bullets: list):
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title_text
            _bold_title(sl)
            tf = sl.placeholders[1].text_frame
            tf.clear()
            first = True
            for line in bullets:
                line = str(line)[:300]
                if not line:
                    continue
                if first:
                    p = tf.paragraphs[0]; first = False
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                if line.startswith("  ") or line.startswith("\t"):
                    p.level = 1
            return sl

        # ── 1. Title ──────────────────────────────────────────────────────────
        sl1 = prs.slides.add_slide(prs.slide_layouts[0])
        sl1.shapes.title.text = c["title"]
        sl1.placeholders[1].text = (
            f"{c['authors']}\n"
            f"{c['year']}  ·  {c['source']}\n\n"
            f"Submitted for Academic Presentation"
        )
        _bold_title(sl1)

        # ── 2. Agenda ─────────────────────────────────────────────────────────
        add_slide("Agenda", [
            "1.  Introduction",
            "2.  Literature Survey",
            "3.  Problem Statement / Limitations",
            "4.  Objective",
            "5.  Proposed System",
            "6.  Architecture Diagram",
            "7.  Modules",
            "8.  System Requirements",
            "9.  Results",
            "10. Conclusion",
            "11. Future Enhancements",
            "12. References",
        ])

        # ── 3. Introduction ───────────────────────────────────────────────────
        intro = [f"• {s}" for s in c["intro"]]
        if not intro:
            intro = [f"• {c['abstract'][:300]}"]
        add_slide("Introduction", intro)

        # ── 4. Literature Survey ──────────────────────────────────────────────
        lit = [
            f"• Existing research in {m[0] if m else 'this domain'} provides foundational methods",
            f"• Prior work on {m[1] if len(m)>1 else 'related topics'} highlights key challenges",
        ]
        if c["problems"]:
            lit.append(f"• {c['problems'][0]}")
        lit += [
            f"• Gap identified in literature on {m[2] if len(m)>2 else 'the proposed approach'}",
            f"• This study extends prior work on {', '.join(m[:3]) if m else 'the topic'}",
        ]
        add_slide("Literature Survey", lit)

        # ── 5. Problem Statement / Limitations ───────────────────────────────
        prob = [f"• {s}" for s in c["problems"]]
        if not prob:
            prob = [
                f"• Existing {m[0] if m else 'approaches'} exhibit performance limitations",
                "• Solutions do not generalize well to real-world conditions",
                "• Lack of efficiency and scalability in prior methods",
            ]
        add_slide("Problem Statement / Limitations", prob)

        # ── 6. Objective ──────────────────────────────────────────────────────
        obj = []
        if c["tldr"]:
            obj.append(f"• {c['tldr']}")
        obj += [f"• {s}" for s in c["objectives"][:3]]
        if not obj:
            obj = [
                f"• To design a novel solution for {c['title'][:70]}",
                "• To evaluate performance against state-of-the-art baselines",
                "• To demonstrate practical applicability on real datasets",
            ]
        add_slide("Objective", obj)

        # ── 7. Proposed System ────────────────────────────────────────────────
        meth = [f"• {s}" for s in c["method"][:4]]
        if not meth:
            meth = [
                f"• Proposed: a novel {m[0] if m else 'system'}-based approach",
                "• End-to-end trainable framework",
                "• Addresses identified limitations through innovative design",
            ]
        add_slide("Proposed System", meth)

        # ── 8. Architecture Diagram ───────────────────────────────────────────
        add_slide("Architecture Diagram", [
            "High-Level System Architecture:",
            "",
            f"  [ Input Data / Dataset ]",
            f"         ↓",
            f"  [ {m[0].title() if m else 'Pre-processing'} Module ]",
            f"         ↓",
            f"  [ {m[1].title() if len(m)>1 else 'Core'} Engine  ←→  {m[2].title() if len(m)>2 else 'Analysis'} Layer ]",
            f"         ↓",
            f"  [ {m[3].title() if len(m)>3 else 'Output'} / Results ]",
            "",
            "  * Refer to paper for the complete architecture diagram.",
        ])

        # ── 9. Modules ────────────────────────────────────────────────────────
        mods = [f"• {w.replace('-', ' ').title()} Module" for w in m[:7]]
        if not mods:
            mods = ["• Input Processing Module", "• Core Analysis Module", "• Output Generation Module"]
        add_slide("Modules", mods)

        # ── 10. System Requirements ───────────────────────────────────────────
        add_slide("System Requirements", [
            "Hardware Requirements:",
            "  • CPU: Intel Core i5 / AMD Ryzen 5 or higher",
            "  • RAM: 8 GB minimum  (16 GB recommended)",
            "  • Storage: 10 GB free disk space",
            "  • GPU: NVIDIA GPU with CUDA support (for DL models)",
            "",
            "Software Requirements:",
            "  • OS: Windows 10 / Ubuntu 20.04 / macOS 12+",
            "  • Python 3.8+  |  pip package manager",
            f"  • Key libraries: {', '.join(m[:5]) if m else 'standard ML/DL stack'}",
        ])

        # ── 11. Results ───────────────────────────────────────────────────────
        res = [f"• {s}" for s in c["results"][:4]]
        if not res:
            res = [
                "• Proposed method evaluated on standard benchmark datasets",
                "• Achieves state-of-the-art performance on key evaluation metrics",
                "• Significant improvement in accuracy / efficiency over baselines",
            ]
        add_slide("Results", res)

        # ── 12. Conclusion ────────────────────────────────────────────────────
        conc = []
        if c["tldr"]:
            conc.append(f"• {c['tldr']}")
        conc += [f"• {kp}" for kp in c["keypoints"][:3]]
        if not conc:
            conc = [
                "• Successfully addressed the identified research challenges",
                "• Proposed approach demonstrates superior performance",
                "• Work contributes meaningfully to the research community",
            ]
        add_slide("Conclusion", conc)

        # ── 13. Future Enhancements ───────────────────────────────────────────
        fut = [f"• {s}" for s in c["future"][:3]]
        if not fut:
            fut = [
                "• Extend to larger and more diverse real-world datasets",
                "• Optimize for real-time deployment and edge computing",
                "• Explore cross-domain transfer learning capabilities",
                "• Investigate integration with complementary methods",
            ]
        add_slide("Future Enhancements", fut)

        # ── 14. References ────────────────────────────────────────────────────
        add_slide("References", [
            f"[1] {c['authors']} ({c['year']}).",
            f"    {c['title'][:100]}.",
            f"    {c['source']}.",
            f"    URL: {c['url']}",
            "",
            "For the complete reference list, refer to the original paper.",
        ])

        # ── Save ──────────────────────────────────────────────────────────────
        safe = re.sub(r"[^\w\s-]", "", paper["title"])[:40].strip().replace(" ", "_")
        path = str(PPTX / f"{safe}.pptx")
        prs.save(path)
        return path

    except Exception as e:
        console.print(f"[red]PPTX error: {e}[/red]")
        import traceback; traceback.print_exc()
        return None


# ══════════════════════════════════════════════════════════════════════════════
# KNOWLEDGE GRAPH
# ══════════════════════════════════════════════════════════════════════════════

STOP = {"the","a","an","of","in","for","and","or","to","is","are","with","that",
        "this","we","our","on","at","by","from","as","be","it","its","also","but",
        "not","have","has","was","were","using","used","based","which","paper",
        "show","can","when","their","they","been","more","than","these","model"}

def build_graph(papers: List[Dict]):
    concept_map: Dict[str, List[str]] = {}
    paper_concepts = {}

    for p in papers:
        text  = (p["title"] + " " + p.get("abstract","")).lower()
        words = re.findall(r"\b[a-z]{4,}\b", text)
        freq  = {}
        for w in words:
            if w not in STOP: freq[w] = freq.get(w,0)+1
        top = sorted(freq, key=freq.get, reverse=True)[:8]
        paper_concepts[p["id"]] = top
        for c in top:
            concept_map.setdefault(c, []).append(p["id"])

    shared = {k: v for k, v in concept_map.items() if len(v) > 1}
    return paper_concepts, shared


# ══════════════════════════════════════════════════════════════════════════════
# GAP DISCOVERY
# ══════════════════════════════════════════════════════════════════════════════

def discover_gaps(papers: List[Dict]) -> List[Dict]:
    gaps = []
    all_c: Dict[str, List[str]] = {}
    for p in papers:
        text  = (p["title"] + " " + p.get("abstract","")).lower()
        words = re.findall(r"\b[a-z]{4,}\b", text)
        for w in set(words):
            if w not in STOP and len(w) > 5:
                all_c.setdefault(w, []).append(p["id"])

    single = sorted([(c,pids) for c,pids in all_c.items() if len(pids)==1], key=lambda x:len(x[0]), reverse=True)
    for c, pids in single[:4]:
        gaps.append({"type":"Underexplored","title":f"'{c}' studied by only 1 paper","confidence":72,"impact":"Medium"})

    years = [(p.get("year",0), p["title"]) for p in papers if p.get("year",0) > 0]
    if years:
        min_y, max_y = min(y[0] for y in years), max(y[0] for y in years)
        if max_y - min_y > 3:
            for yr, title in sorted(years)[:2]:
                gaps.append({"type":"Temporal Gap","title":f"'{title[:45]}' ({yr}) — ripe for modern re-evaluation","confidence":68,"impact":"High"})

    for c, pids in all_c.items():
        if 2 <= len(pids) <= 5 and len(c) > 6:
            src = set()
            for pid in pids:
                paper = next((p for p in papers if p["id"]==pid), None)
                if paper: src.add(paper.get("source",""))
            if len(src) > 1:
                gaps.append({"type":"Cross-Domain","title":f"'{c}' bridges multiple research communities","confidence":81,"impact":"Breakthrough"})
            if len(gaps) >= 12: break

    for p in papers:
        abst = p.get("abstract","").lower()
        if "baseline" in abst and "comparison" not in abst:
            gaps.append({"type":"Methodological","title":f"Missing comparison: '{p['title'][:45]}'","confidence":75,"impact":"Medium"})

    return gaps[:12]


# ══════════════════════════════════════════════════════════════════════════════
# DISPLAY HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _source_tag(source: str) -> str:
    if "arXiv" in source: return "[bold cyan]arXiv[/bold cyan]"
    return "[bold magenta]S2[/bold magenta]"

def show_papers(papers: List[Dict], title: str = ""):
    if not papers:
        console.print("[yellow]No papers.[/yellow]")
        return

    t = Table(
        box=box.ROUNDED,
        header_style="bold bright_white on #1a1a2e",
        show_lines=False,
        padding=(0, 1),
        title=f"[bold]{title}[/bold]" if title else None,
        title_style="bold cyan",
    )
    t.add_column("#",        width=3,  style="dim", justify="right")
    t.add_column("Source",   width=7,  justify="center")
    t.add_column("Title",    width=52)
    t.add_column("Authors",  width=22, style="dim")
    t.add_column("Year",     width=6,  justify="center", style="green")
    t.add_column("Cit.",     width=6,  justify="right",  style="yellow")

    for i, p in enumerate(papers, 1):
        auth = ", ".join(p.get("authors",[])[:2])
        if len(p.get("authors",[])) > 2: auth += " …"
        t.add_row(
            str(i),
            _source_tag(p.get("source","")),
            (p["title"][:50] + "…") if len(p["title"]) > 50 else p["title"],
            auth[:20],
            str(p.get("year") or "—"),
            str(p.get("citations") or "0"),
        )
    console.print(t)

def show_summary_card(paper: Dict, summary: Dict):
    console.print()
    console.print(Panel(
        f"[bold white]{paper['title']}[/bold white]\n"
        f"[dim]{', '.join(paper.get('authors',[])[:3])} · {paper.get('year','')} · {paper.get('source','')}[/dim]",
        border_style="cyan", padding=(0,1)
    ))

    # TL;DR
    console.print(Panel(
        f"[bold yellow]TL;DR[/bold yellow]\n[white]{summary['tldr']}[/white]",
        border_style="yellow", padding=(0,1)
    ))

    # Key Points
    if summary.get("keypoints"):
        kp = "\n".join(f"  [cyan]•[/cyan] {k}" for k in summary["keypoints"])
        console.print(Panel(
            f"[bold green]Key Contributions[/bold green]\n{kp}",
            border_style="green", padding=(0,1)
        ))

    # Full Summary
    if summary.get("full") and summary["full"] != paper.get("abstract",""):
        console.print(Panel(
            f"[bold blue]Full Summary[/bold blue]\n[dim white]{summary['full']}[/dim white]",
            border_style="blue", padding=(0,1)
        ))
    console.print()

def show_help():
    console.print()
    console.print(Panel(
        "[bold cyan]ScholarGenie Commands[/bold cyan]\n\n"
        "[bold yellow]search[/bold yellow] [green]<topic>[/green] [dim][--count N][/dim]\n"
        "   Search arXiv + Semantic Scholar for papers\n\n"
        "[bold yellow]summarize[/bold yellow] [green]<N>[/green] [dim](or 'all')[/dim]\n"
        "   Summarize paper N from last search / library\n\n"
        "[bold yellow]present[/bold yellow] [green]<N>[/green]\n"
        "   Generate PowerPoint for paper N\n\n"
        "[bold yellow]open[/bold yellow] [green]<N>[/green]\n"
        "   Open paper PDF/URL in browser\n\n"
        "[bold yellow]library[/bold yellow]\n"
        "   View all saved papers\n\n"
        "[bold yellow]graph[/bold yellow]\n"
        "   Show knowledge graph of saved papers\n\n"
        "[bold yellow]gaps[/bold yellow]\n"
        "   Discover research gaps in saved papers\n\n"
        "[bold yellow]workflow[/bold yellow] [green]<topic>[/green]\n"
        "   Full pipeline: Search → Summarize → Present → Graph → Gaps\n\n"
        "[bold yellow]clear[/bold yellow]  [dim]Clear screen[/dim]\n"
        "[bold yellow]exit[/bold yellow]   [dim]Quit ScholarGenie[/dim]",
        title="[bold]Help[/bold]",
        border_style="dim",
        padding=(0, 2),
    ))
    console.print()

def show_banner():
    console.clear()
    art = (
        " ____       _           _           ____            _      \n"
        "/ ___|  ___| |__   ___ | | __ _ _ _/ ___| ___ _ __ (_) ___ \n"
        "\\___ \\ / __| '_ \\ / _ \\| |/ _` | '__\\___ \\/ _ \\ '_ \\| |/ _ \\\n"
        " ___) | (__| | | | (_) | | (_| | |   ___) |  __/ | | | |  __/\n"
        "|____/ \\___|_| |_|\\___/|_|\\__,_|_|  |____/ \\___|_| |_|_|\\___|\n"
    )
    console.print(Panel(
        f"[bold cyan]{art}[/bold cyan]\n"
        "[bold white]  AI Research Intelligence Platform[/bold white]  "
        "[dim]· arXiv · Semantic Scholar · Summaries · Slides · Gaps[/dim]",
        border_style="cyan",
        padding=(0, 2),
    ))
    saved = db.count()
    console.print(
        f"  [dim]Library: [cyan]{saved}[/cyan] papers saved  ·  "
        f"Data: [cyan]{DATA}[/cyan]  ·  "
        f"Type [bold yellow]help[/bold yellow] for commands[/dim]\n"
    )


# ══════════════════════════════════════════════════════════════════════════════
# COMMAND HANDLERS
# ══════════════════════════════════════════════════════════════════════════════

def cmd_search(args: str):
    """search <topic> [--count N]"""
    # Parse --count flag
    n = 10
    m = re.search(r"--count\s+(\d+)", args)
    if m:
        n = int(m.group(1))
        args = args[:m.start()].strip()
    query = args.strip()
    if not query:
        console.print("[red]Usage: search <topic> [--count N][/red]")
        return

    console.print(f"\n  [dim]Searching for:[/dim] [bold]{query}[/bold]  [dim]({n} results)[/dim]\n")

    with Progress(
        SpinnerColumn(style="cyan"),
        TextColumn("[cyan]{task.description}"),
        TimeElapsedColumn(),
        console=console,
        transient=True,
    ) as prog:
        t1 = prog.add_task("Searching arXiv...", total=None)
        ax = _arxiv(query, n)
        prog.update(t1, description=f"[green]arXiv: {len(ax)} found[/green]")

        t2 = prog.add_task("Searching Semantic Scholar...", total=None)
        ss = _semantic(query, n)
        prog.update(t2, description=f"[green]Semantic Scholar: {len(ss)} found[/green]")

    all_p = ax + ss
    seen, unique = set(), []
    for p in all_p:
        key = re.sub(r"\W+", "", p["title"].lower())[:60]
        if key not in seen and len(key) > 4:
            seen.add(key); unique.append(p)
    results = unique[:n]

    if not results:
        console.print("[red]No papers found. Check your internet connection.[/red]\n")
        return

    for p in results:
        db.upsert(p)

    SESSION["results"] = results
    SESSION["query"]   = query

    show_papers(results, title=f'"{query}" — {len(results)} papers')
    console.print(
        f"  [dim]Saved {len(results)} papers to library.  "
        f"Use [bold yellow]summarize N[/bold yellow] or [bold yellow]present N[/bold yellow][/dim]\n"
    )


def cmd_summarize(args: str):
    """summarize <N> | all"""
    arg = args.strip().lower()
    pool = SESSION["results"] or db.all_papers()
    if not pool:
        console.print("[yellow]No papers. Run 'search <topic>' first.[/yellow]\n")
        return

    targets = []
    if arg == "all":
        targets = pool
    elif arg.isdigit():
        idx = int(arg)
        if 1 <= idx <= len(pool):
            targets = [pool[idx-1]]
        else:
            console.print(f"[red]Paper {idx} not found. Range is 1–{len(pool)}.[/red]\n")
            return
    else:
        console.print("[red]Usage: summarize <N>  or  summarize all[/red]\n")
        return

    for paper in targets:
        console.print(f"\n  [cyan]Summarizing:[/cyan] [white]{paper['title'][:65]}...[/white]")
        with console.status("  Generating summary...", spinner="dots"):
            s = summarize(paper)
        db.save_summary(paper["id"], s)
        show_summary_card(paper, s)
        if len(targets) > 1: time.sleep(0.1)


def cmd_present(args: str):
    """present <N>"""
    arg = args.strip()
    pool = SESSION["results"] or db.all_papers()
    if not pool:
        console.print("[yellow]No papers. Run 'search <topic>' first.[/yellow]\n")
        return
    if not arg.isdigit():
        console.print("[red]Usage: present <N>[/red]\n")
        return
    idx = int(arg)
    if not (1 <= idx <= len(pool)):
        console.print(f"[red]Paper {idx} not found. Range is 1–{len(pool)}.[/red]\n")
        return

    paper = pool[idx-1]
    console.print(f"\n  [cyan]Building slides for:[/cyan] [white]{paper['title'][:65]}[/white]")

    with console.status("  Summarizing...", spinner="dots"):
        summary = db.get_summary(paper["id"]) or summarize(paper)

    with console.status("  Generating PowerPoint...", spinner="dots"):
        path = make_pptx(paper, summary)

    if path:
        console.print(Panel(
            f"[bold green]Presentation ready![/bold green]\n\n"
            f"  [cyan]File:[/cyan]   {path}\n\n"
            f"  [cyan]Slides (14):[/cyan]\n"
            f"  Title · Agenda · Introduction · Literature Survey ·\n"
            f"  Problem Statement · Objective · Proposed System ·\n"
            f"  Architecture · Modules · System Requirements ·\n"
            f"  Results · Conclusion · Future Enhancements · References",
            border_style="green", padding=(0,1)
        ))
        try:
            os.startfile(path) if sys.platform == "win32" else None
        except Exception:
            pass
    else:
        console.print("[red]Failed to create presentation. Is python-pptx installed?[/red]\n")


def cmd_open(args: str):
    """open <N>"""
    arg = args.strip()
    pool = SESSION["results"] or db.all_papers()
    if not pool:
        console.print("[yellow]No papers in session. Run 'search' first.[/yellow]\n")
        return
    if not arg.isdigit():
        console.print("[red]Usage: open <N>[/red]\n")
        return
    idx = int(arg)
    if not (1 <= idx <= len(pool)):
        console.print(f"[red]Out of range. Use 1–{len(pool)}.[/red]\n")
        return
    url = pool[idx-1].get("url", "")
    if url:
        webbrowser.open(url)
        console.print(f"  [green]Opened:[/green] {url}\n")
    else:
        console.print("[yellow]No URL for this paper.[/yellow]\n")


def cmd_library():
    papers = db.all_papers()
    if not papers:
        console.print("[yellow]Library is empty. Run 'search <topic>' to add papers.[/yellow]\n")
        return
    show_papers(papers, title=f"Library — {len(papers)} papers")
    SESSION["results"] = papers
    console.print(
        f"  [dim]Library loaded as active pool. "
        f"You can now use summarize/present/open by number.[/dim]\n"
    )


def cmd_graph():
    pool = SESSION["results"] or db.all_papers()
    if not pool:
        console.print("[yellow]No papers. Run 'search' or 'library' first.[/yellow]\n")
        return

    console.print(f"\n  [cyan]Building knowledge graph from {len(pool)} papers...[/cyan]\n")
    paper_concepts, shared = build_graph(pool)

    # Paper → concepts table
    t = Table(title="[bold]Paper → Key Concepts[/bold]", box=box.ROUNDED,
              header_style="bold bright_white", padding=(0,1))
    t.add_column("#",        width=3,  justify="right", style="dim")
    t.add_column("Paper",    width=46, style="cyan")
    t.add_column("Concepts", width=52, style="green")

    for i, p in enumerate(pool, 1):
        concepts = paper_concepts.get(p["id"], [])
        t.add_row(str(i), p["title"][:44], "  •  ".join(concepts[:6]))
    console.print(t)

    # Shared concepts
    if shared:
        console.print(f"\n  [bold]Shared Concepts[/bold] [dim]({len(shared)} found)[/dim]\n")
        rows = sorted(shared.items(), key=lambda x: len(x[1]), reverse=True)[:15]
        st = Table(box=box.SIMPLE, padding=(0,1), show_header=False)
        st.add_column("Concept", style="bold yellow", width=22)
        st.add_column("Papers",  style="dim white")
        for concept, pids in rows:
            titles = []
            for pid in pids[:3]:
                p = next((x for x in pool if x["id"]==pid), None)
                if p: titles.append(p["title"][:30])
            st.add_row(concept, " ↔ ".join(titles))
        console.print(st)
    else:
        console.print("  [dim]No shared concepts found — try searching more papers on the same topic.[/dim]")
    console.print()


def cmd_gaps():
    pool = SESSION["results"] or db.all_papers()
    if not pool:
        console.print("[yellow]No papers. Run 'search' or 'library' first.[/yellow]\n")
        return

    console.print(f"\n  [cyan]Analyzing {len(pool)} papers for research gaps...[/cyan]\n")
    with console.status("  Running gap analysis...", spinner="dots"):
        gaps = discover_gaps(pool)

    if not gaps:
        console.print("  [dim]No gaps detected in current corpus.[/dim]\n")
        return

    impact_colors = {"Breakthrough": "bold green", "High": "bold orange3",
                     "Medium": "yellow", "Low": "dim"}
    type_colors   = {"Cross-Domain": "green", "Temporal Gap": "blue",
                     "Methodological": "yellow", "Underexplored": "magenta"}

    console.print(f"  [bold]Found {len(gaps)} research gaps[/bold]\n")
    for i, gap in enumerate(gaps, 1):
        impact_col = impact_colors.get(gap["impact"], "white")
        type_col   = type_colors.get(gap["type"], "cyan")
        console.print(
            f"  [{type_col}]{i:2}. {gap['type']:<16}[/{type_col}]  "
            f"[{impact_col}]{gap['impact']:<12}[/{impact_col}]  "
            f"[dim]{gap['confidence']}%[/dim]  "
            f"[white]{gap['title']}[/white]"
        )
    console.print()


def cmd_workflow(args: str):
    """Full pipeline: Search → Summarize → Present → Graph → Gaps"""
    m = re.search(r"--count\s+(\d+)", args)
    n = int(m.group(1)) if m else 5
    if m: args = args[:m.start()].strip()
    query = args.strip()
    if not query:
        console.print("[red]Usage: workflow <topic> [--count N][/red]\n")
        return

    console.print(Panel(
        f"[bold cyan]Full Research Workflow[/bold cyan]\n"
        f"[dim]Search → Summarize All → Generate Slides → Knowledge Graph → Research Gaps[/dim]\n"
        f"[white]Topic: {query}  ·  Papers: {n}[/white]",
        border_style="cyan", padding=(0,1)
    ))

    # ── 1. Search ────────────────────────────────────────────────────────────
    console.print("\n  [bold]1/5  Searching...[/bold]")
    with Progress(SpinnerColumn(style="cyan"), TextColumn("[cyan]{task.description}"),
                  console=console, transient=True) as prog:
        t = prog.add_task("Searching arXiv + Semantic Scholar...", total=None)
        ax = _arxiv(query, n)
        ss = _semantic(query, n)

    all_p = ax + ss
    seen, unique = set(), []
    for p in all_p:
        key = re.sub(r"\W+","",p["title"].lower())[:60]
        if key not in seen and len(key)>4: seen.add(key); unique.append(p)
    papers = unique[:n]

    if not papers:
        console.print("[red]No papers found.[/red]\n")
        return

    for p in papers: db.upsert(p)
    SESSION["results"] = papers; SESSION["query"] = query
    show_papers(papers, title=f'"{query}" — {len(papers)} papers')

    # ── 2. Summarize all ──────────────────────────────────────────────────────
    console.print(f"\n  [bold]2/5  Summarizing {len(papers)} papers...[/bold]")
    summaries = {}
    for i, p in enumerate(papers, 1):
        console.print(f"  [{i}/{len(papers)}] [dim]{p['title'][:58]}[/dim]")
        with console.status("       Summarizing...", spinner="dots"):
            s = summarize(p)
        db.save_summary(p["id"], s)
        summaries[p["id"]] = s
    console.print(f"  [green]All {len(papers)} papers summarized.[/green]")

    # ── 3. Presentation for top paper ─────────────────────────────────────────
    console.print(f"\n  [bold]3/5  Generating presentation for top paper...[/bold]")
    top   = papers[0]
    top_s = summaries[top["id"]]
    with console.status("       Building PowerPoint...", spinner="dots"):
        ppath = make_pptx(top, top_s)
    if ppath:
        console.print(f"  [green]Saved:[/green] {ppath}")
    else:
        console.print("  [yellow]Presentation skipped (python-pptx missing).[/yellow]")

    # ── 4. Knowledge Graph ────────────────────────────────────────────────────
    console.print(f"\n  [bold]4/5  Building knowledge graph...[/bold]")
    paper_concepts, shared = build_graph(papers)
    console.print(f"  [green]Graph: {len(shared)} shared concepts across {len(papers)} papers.[/green]")
    for c, pids in list(shared.items())[:6]:
        titles = [p["title"][:28] for p in papers if p["id"] in pids]
        console.print(f"  [yellow]{c}[/yellow]: {' ↔ '.join(titles)}")

    # ── 5. Gaps ────────────────────────────────────────────────────────────────
    console.print(f"\n  [bold]5/5  Discovering research gaps...[/bold]")
    gaps = discover_gaps(papers)
    console.print(f"  [green]Found {len(gaps)} gaps.[/green]")
    for g in gaps[:5]:
        console.print(f"  [cyan]•[/cyan] [{g['type']}] {g['title']}")

    # ── Report ─────────────────────────────────────────────────────────────────
    console.print()
    console.print(Panel(
        f"[bold green]Workflow Complete![/bold green]\n\n"
        f"  [cyan]Papers found:[/cyan]      {len(papers)}\n"
        f"  [cyan]Summaries:[/cyan]         {len(papers)}\n"
        f"  [cyan]Presentation:[/cyan]      {ppath or 'N/A'}\n"
        f"  [cyan]Shared concepts:[/cyan]   {len(shared)}\n"
        f"  [cyan]Research gaps:[/cyan]     {len(gaps)}\n\n"
        f"  [dim]Use [bold yellow]summarize N[/bold yellow] for detailed view  ·  [bold yellow]open N[/bold yellow] to open PDF[/dim]",
        border_style="green", padding=(0,1)
    ))
    console.print()
    # Show top paper summary
    show_summary_card(top, top_s)


# ══════════════════════════════════════════════════════════════════════════════
# REPL MAIN LOOP
# ══════════════════════════════════════════════════════════════════════════════

PROMPT = "[bold cyan]sg[/bold cyan] [dim white]›[/dim white] "

def run():
    show_banner()
    console.print("  Type [bold yellow]help[/bold yellow] for commands  or  [bold yellow]search <topic>[/bold yellow] to get started\n")

    while True:
        try:
            console.print(PROMPT, end="")
            raw = input().strip()
        except (EOFError, KeyboardInterrupt):
            console.print("\n[cyan]Goodbye![/cyan]\n")
            break

        if not raw:
            continue

        parts = raw.split(None, 1)
        cmd   = parts[0].lower()
        args  = parts[1] if len(parts) > 1 else ""

        if cmd in ("exit", "quit", "q"):
            console.print("\n[cyan]Goodbye![/cyan]\n")
            break
        elif cmd in ("help", "h", "?"):
            show_help()
        elif cmd == "clear":
            show_banner()
        elif cmd == "search":
            cmd_search(args)
        elif cmd == "summarize":
            cmd_summarize(args)
        elif cmd == "present":
            cmd_present(args)
        elif cmd == "open":
            cmd_open(args)
        elif cmd == "library":
            cmd_library()
        elif cmd == "graph":
            cmd_graph()
        elif cmd == "gaps":
            cmd_gaps()
        elif cmd == "workflow":
            cmd_workflow(args)
        else:
            console.print(
                f"  [red]Unknown command:[/red] [bold]{cmd}[/bold]  "
                f"[dim]Type [bold yellow]help[/bold yellow] for available commands.[/dim]\n"
            )


if __name__ == "__main__":
    try:
        run()
    except Exception as e:
        console.print(f"\n[red]Fatal error: {e}[/red]")
        import traceback
        traceback.print_exc()
