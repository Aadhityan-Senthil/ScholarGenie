"""Presentation and report generation agent."""

import os
import logging
from typing import Optional, List
from datetime import datetime
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import markdown

from backend.utils.metadata import PaperMetadata, Summary, ExtractedData

logger = logging.getLogger(__name__)


class PresenterAgent:
    """Agent for generating presentations and reports."""

    def __init__(self, config_path: str = "config.yaml"):
        """Initialize presenter agent.

        Args:
            config_path: Path to configuration file
        """
        with open(config_path, 'r') as f:
            self.config = yaml.safe_load(f)

        self.pres_config = self.config["presentation"]
        self.report_config = self.config["report"]

    def generate_pptx(
        self,
        paper: PaperMetadata,
        summary: Summary,
        extracted_data: Optional[ExtractedData] = None,
        output_path: Optional[str] = None
    ) -> str:
        """Generate PowerPoint presentation.

        Args:
            paper: Paper metadata
            summary: Paper summary
            extracted_data: Optional extracted data
            output_path: Output file path

        Returns:
            Path to generated PPTX file
        """
        logger.info(f"Generating presentation for: {paper.title}")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Add slides
        self._add_title_slide(prs, paper)
        self._add_overview_slide(prs, summary)

        if extracted_data:
            self._add_motivation_slide(prs, extracted_data)
            self._add_methods_slide(prs, extracted_data, summary)
            self._add_results_slide(prs, extracted_data, summary)
            self._add_limitations_slide(prs, extracted_data, summary)
        else:
            # Fallback without extracted data
            if summary.methods:
                self._add_methods_slide_simple(prs, summary)
            if summary.results:
                self._add_results_slide_simple(prs, summary)

        self._add_keypoints_slide(prs, summary)

        if paper.references:
            self._add_references_slide(prs, paper)

        # Save presentation
        if not output_path:
            output_dir = self.config["storage"]["output_dir"]
            os.makedirs(output_dir, exist_ok=True)
            filename = f"{paper.paper_id}_presentation.pptx"
            output_path = os.path.join(output_dir, filename)

        prs.save(output_path)
        logger.info(f"Presentation saved to: {output_path}")

        return output_path

    def _add_title_slide(self, prs: Presentation, paper: PaperMetadata):
        """Add title slide.

        Args:
            prs: Presentation object
            paper: Paper metadata
        """
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = paper.title

        # Add authors and venue
        author_names = [a.name for a in paper.authors[:5]]  # Limit to 5
        if len(paper.authors) > 5:
            author_names.append("et al.")

        subtitle_text = ", ".join(author_names)
        if paper.venue:
            subtitle_text += f"\n{paper.venue}"
        if paper.year:
            subtitle_text += f"\n{paper.year}"

        subtitle.text = subtitle_text

    def _add_overview_slide(self, prs: Presentation, summary: Summary):
        """Add overview slide with TL;DR and short summary.

        Args:
            prs: Presentation object
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content

        title = slide.shapes.title
        title.text = "Overview"

        # Add text box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # TL;DR
        p = tf.paragraphs[0]
        p.text = "TL;DR"
        p.font.bold = True
        p.font.size = Pt(16)

        # Add TL;DR content
        p = tf.add_paragraph()
        p.text = summary.tldr
        p.font.size = Pt(14)
        p.level = 1

        # Short summary
        p = tf.add_paragraph()
        p.text = "\nSummary"
        p.font.bold = True
        p.font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = summary.short_summary
        p.font.size = Pt(12)
        p.level = 1

    def _add_motivation_slide(self, prs: Presentation, extracted: ExtractedData):
        """Add motivation/research question slide.

        Args:
            prs: Presentation object
            extracted: Extracted data
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Research Question & Motivation"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        if extracted.research_question:
            p = tf.paragraphs[0]
            p.text = "Research Question"
            p.font.bold = True
            p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = extracted.research_question
            p.font.size = Pt(12)
            p.level = 1

        if extracted.hypothesis:
            p = tf.add_paragraph()
            p.text = "\nHypothesis"
            p.font.bold = True
            p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = extracted.hypothesis
            p.font.size = Pt(12)
            p.level = 1

    def _add_methods_slide(
        self,
        prs: Presentation,
        extracted: ExtractedData,
        summary: Summary
    ):
        """Add methods slide.

        Args:
            prs: Presentation object
            extracted: Extracted data
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Methodology"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Methods summary
        if summary.methods:
            p = tf.paragraphs[0]
            p.text = summary.methods
            p.font.size = Pt(12)

        # Key methods
        if extracted.methods:
            p = tf.add_paragraph()
            p.text = "\nKey Techniques:"
            p.font.bold = True
            p.font.size = Pt(14)

            for method in extracted.methods[:8]:
                p = tf.add_paragraph()
                p.text = method
                p.font.size = Pt(11)
                p.level = 1

        # Datasets
        if extracted.datasets:
            p = tf.add_paragraph()
            p.text = "\nDatasets: " + ", ".join(extracted.datasets[:5])
            p.font.size = Pt(11)

    def _add_methods_slide_simple(self, prs: Presentation, summary: Summary):
        """Add simple methods slide without extracted data.

        Args:
            prs: Presentation object
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Methodology"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = summary.methods if summary.methods else "See paper for details."
        p.font.size = Pt(12)

    def _add_results_slide(
        self,
        prs: Presentation,
        extracted: ExtractedData,
        summary: Summary
    ):
        """Add results slide.

        Args:
            prs: Presentation object
            extracted: Extracted data
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Results & Findings"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Results summary
        if summary.results:
            p = tf.paragraphs[0]
            p.text = summary.results
            p.font.size = Pt(12)

        # Key findings
        if extracted.key_findings:
            p = tf.add_paragraph()
            p.text = "\nKey Findings:"
            p.font.bold = True
            p.font.size = Pt(14)

            for finding in extracted.key_findings[:6]:
                p = tf.add_paragraph()
                p.text = finding
                p.font.size = Pt(11)
                p.level = 1

        # Metrics
        if extracted.metrics:
            p = tf.add_paragraph()
            p.text = "\nPerformance Metrics:"
            p.font.bold = True
            p.font.size = Pt(14)

            for metric, value in list(extracted.metrics.items())[:5]:
                p = tf.add_paragraph()
                p.text = f"{metric}: {value}"
                p.font.size = Pt(11)
                p.level = 1

    def _add_results_slide_simple(self, prs: Presentation, summary: Summary):
        """Add simple results slide.

        Args:
            prs: Presentation object
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Results"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = summary.results if summary.results else "See paper for details."
        p.font.size = Pt(12)

    def _add_keypoints_slide(self, prs: Presentation, summary: Summary):
        """Add key points slide.

        Args:
            prs: Presentation object
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Key Takeaways"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(summary.keypoints[:10], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i}. {point}"
            p.font.size = Pt(12)
            p.space_after = Pt(8)

    def _add_limitations_slide(
        self,
        prs: Presentation,
        extracted: ExtractedData,
        summary: Summary
    ):
        """Add limitations and future work slide.

        Args:
            prs: Presentation object
            extracted: Extracted data
            summary: Summary object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Limitations & Future Work"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Limitations
        if extracted.limitations:
            p = tf.paragraphs[0]
            p.text = "Limitations:"
            p.font.bold = True
            p.font.size = Pt(14)

            for lim in extracted.limitations[:4]:
                p = tf.add_paragraph()
                p.text = lim
                p.font.size = Pt(11)
                p.level = 1

        # Future work
        if extracted.future_work:
            p = tf.add_paragraph()
            p.text = "\nFuture Directions:"
            p.font.bold = True
            p.font.size = Pt(14)

            for fw in extracted.future_work[:4]:
                p = tf.add_paragraph()
                p.text = fw
                p.font.size = Pt(11)
                p.level = 1

    def _add_references_slide(self, prs: Presentation, paper: PaperMetadata):
        """Add references slide.

        Args:
            prs: Presentation object
            paper: Paper metadata
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Key References"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        max_refs = self.pres_config.get("max_references", 10)

        for i, ref in enumerate(paper.references[:max_refs], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()

            ref_text = f"{i}. {ref.title}"
            if ref.authors:
                ref_text += f" ({', '.join(ref.authors[:2])})"
            if ref.year:
                ref_text += f", {ref.year}"

            p.text = ref_text
            p.font.size = Pt(9)
            p.space_after = Pt(6)

    def generate_markdown_report(
        self,
        paper: PaperMetadata,
        summary: Summary,
        extracted_data: Optional[ExtractedData] = None,
        output_path: Optional[str] = None
    ) -> str:
        """Generate Markdown report.

        Args:
            paper: Paper metadata
            summary: Paper summary
            extracted_data: Optional extracted data
            output_path: Output file path

        Returns:
            Path to generated Markdown file
        """
        logger.info(f"Generating Markdown report for: {paper.title}")

        # Build Markdown content
        md_content = self._build_markdown_content(paper, summary, extracted_data)

        # Save file
        if not output_path:
            output_dir = self.config["storage"]["output_dir"]
            os.makedirs(output_dir, exist_ok=True)
            filename = f"{paper.paper_id}_report.md"
            output_path = os.path.join(output_dir, filename)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)

        logger.info(f"Markdown report saved to: {output_path}")

        return output_path

    def _build_markdown_content(
        self,
        paper: PaperMetadata,
        summary: Summary,
        extracted_data: Optional[ExtractedData]
    ) -> str:
        """Build Markdown content for report.

        Args:
            paper: Paper metadata
            summary: Summary object
            extracted_data: Optional extracted data

        Returns:
            Markdown content string
        """
        md = f"# {paper.title}\n\n"

        # Metadata
        if self.report_config.get("include_metadata", True):
            md += "## Metadata\n\n"

            author_names = [a.name for a in paper.authors]
            md += f"**Authors:** {', '.join(author_names)}\n\n"

            if paper.venue:
                md += f"**Venue:** {paper.venue}\n\n"

            if paper.year:
                md += f"**Year:** {paper.year}\n\n"

            if paper.doi:
                md += f"**DOI:** {paper.doi}\n\n"

            if paper.arxiv_id:
                md += f"**arXiv ID:** {paper.arxiv_id}\n\n"

            md += "---\n\n"

        # TL;DR
        md += "## TL;DR\n\n"
        md += f"{summary.tldr}\n\n"

        # Summary
        md += "## Summary\n\n"
        md += f"{summary.full_summary}\n\n"

        # Key Points
        md += "## Key Points\n\n"
        for i, point in enumerate(summary.keypoints, 1):
            md += f"{i}. {point}\n"
        md += "\n"

        # Research Question (if available)
        if extracted_data and extracted_data.research_question:
            md += "## Research Question\n\n"
            md += f"{extracted_data.research_question}\n\n"

        # Methods
        if summary.methods:
            md += "## Methodology\n\n"
            md += f"{summary.methods}\n\n"

            if extracted_data and extracted_data.methods:
                md += "**Key Techniques:**\n"
                for method in extracted_data.methods:
                    md += f"- {method}\n"
                md += "\n"

            if extracted_data and extracted_data.datasets:
                md += "**Datasets:** " + ", ".join(extracted_data.datasets) + "\n\n"

        # Results
        if summary.results:
            md += "## Results\n\n"
            md += f"{summary.results}\n\n"

            if extracted_data and extracted_data.metrics:
                md += "**Performance Metrics:**\n"
                for metric, value in extracted_data.metrics.items():
                    md += f"- {metric}: {value}\n"
                md += "\n"

        # Limitations
        if summary.limitations:
            md += "## Limitations\n\n"
            md += f"{summary.limitations}\n\n"

        # Section Summaries
        if summary.section_summaries:
            md += "## Detailed Section Summaries\n\n"
            for section_title, section_summary in summary.section_summaries.items():
                md += f"### {section_title}\n\n"
                md += f"{section_summary}\n\n"

        # References
        if self.report_config.get("include_references", True) and paper.references:
            md += "## References\n\n"
            for i, ref in enumerate(paper.references[:20], 1):
                ref_text = f"{i}. {ref.title}"
                if ref.authors:
                    ref_text += f" ({', '.join(ref.authors)})"
                if ref.year:
                    ref_text += f", {ref.year}"
                if ref.doi:
                    ref_text += f". DOI: {ref.doi}"

                md += f"{ref_text}\n"

        # Footer
        md += f"\n---\n\n"
        md += f"*Report generated by ScholarGenie on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n"

        return md
