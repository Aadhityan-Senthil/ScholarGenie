#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ScholarGenie API Server
Bridges the CLI engine to the Next.js frontend.
Uses SQLite — no PostgreSQL, Redis, or Docker required.
"""

import os
import sys
import io
import json
import re
import sqlite3
import time
import asyncio
from pathlib import Path
from typing import List, Optional, Dict, Any
import warnings
warnings.filterwarnings("ignore")

# Load .env before anything else
try:
    from dotenv import load_dotenv
    _env_path = Path(__file__).parent / ".env"
    load_dotenv(dotenv_path=_env_path, override=True)
except ImportError:
    pass

# Fix Windows encoding
if sys.platform == "win32":
    try:
        if hasattr(sys.stdout, "buffer"):
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    except Exception:
        pass

# ── FastAPI ────────────────────────────────────────────────────────────────────
from fastapi import FastAPI, HTTPException, BackgroundTasks, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel

# ── Directories ────────────────────────────────────────────────────────────────
BASE_DIR  = Path(__file__).parent
DATA_DIR  = BASE_DIR / "data"
PPTX_DIR  = DATA_DIR / "presentations"
DB_PATH   = DATA_DIR / "scholargenie.db"

for d in [DATA_DIR, PPTX_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ── App ────────────────────────────────────────────────────────────────────────
app = FastAPI(
    title="ScholarGenie API",
    description="AI Research Intelligence Platform",
    version="1.0.0",
    docs_url="/docs",
    redoc_url="/redoc",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],          # frontend on :3000
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ══════════════════════════════════════════════════════════════════════════════
# LIGHTWEIGHT DATABASE  (SQLite)
# ══════════════════════════════════════════════════════════════════════════════

class DB:
    def __init__(self):
        self.conn = sqlite3.connect(str(DB_PATH), check_same_thread=False)
        self._init()

    def _init(self):
        self.conn.executescript("""
            CREATE TABLE IF NOT EXISTS papers (
                id TEXT PRIMARY KEY,
                title TEXT,
                authors TEXT,
                abstract TEXT,
                year INTEGER,
                url TEXT,
                source TEXT,
                citations INTEGER DEFAULT 0,
                created_at TEXT DEFAULT (datetime('now'))
            );
            CREATE TABLE IF NOT EXISTS summaries (
                paper_id TEXT PRIMARY KEY,
                tldr TEXT,
                short_summary TEXT,
                full_summary TEXT,
                keypoints TEXT
            );
            CREATE TABLE IF NOT EXISTS searches (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                query TEXT,
                result_count INTEGER,
                created_at TEXT DEFAULT (datetime('now'))
            );
        """)
        self.conn.commit()

    def save_paper(self, p: dict):
        self.conn.execute(
            "INSERT OR REPLACE INTO papers VALUES (?,?,?,?,?,?,?,?,datetime('now'))",
            (p["id"], p["title"], json.dumps(p.get("authors", [])),
             p.get("abstract", ""), p.get("year", 0),
             p.get("url", ""), p.get("source", ""), p.get("citations", 0))
        )
        self.conn.commit()

    def save_summary(self, pid: str, s: dict):
        self.conn.execute(
            "INSERT OR REPLACE INTO summaries VALUES (?,?,?,?,?)",
            (pid, s.get("tldr", ""), s.get("short", ""),
             s.get("full", ""), json.dumps(s.get("keypoints", [])))
        )
        self.conn.commit()

    def log_search(self, query: str, count: int):
        self.conn.execute(
            "INSERT INTO searches (query, result_count) VALUES (?,?)", (query, count)
        )
        self.conn.commit()

    def get_papers(self) -> List[dict]:
        rows = self.conn.execute(
            "SELECT id,title,authors,abstract,year,url,source,citations,created_at "
            "FROM papers ORDER BY created_at DESC"
        ).fetchall()
        return [
            {"id": r[0], "title": r[1],
             "authors": json.loads(r[2]) if r[2] else [],
             "abstract": r[3], "year": r[4],
             "url": r[5], "source": r[6],
             "citations": r[7], "created_at": r[8]}
            for r in rows
        ]

    def get_paper(self, pid: str) -> Optional[dict]:
        row = self.conn.execute(
            "SELECT id,title,authors,abstract,year,url,source,citations "
            "FROM papers WHERE id=?", (pid,)
        ).fetchone()
        if not row:
            return None
        return {"id": row[0], "title": row[1],
                "authors": json.loads(row[2]) if row[2] else [],
                "abstract": row[3], "year": row[4],
                "url": row[5], "source": row[6], "citations": row[7]}

    def get_summary(self, pid: str) -> Optional[dict]:
        row = self.conn.execute(
            "SELECT tldr,short_summary,full_summary,keypoints FROM summaries WHERE paper_id=?",
            (pid,)
        ).fetchone()
        if not row:
            return None
        return {"tldr": row[0], "short": row[1], "full": row[2],
                "keypoints": json.loads(row[3]) if row[3] else []}

    def delete_paper(self, pid: str):
        self.conn.execute("DELETE FROM papers WHERE id=?", (pid,))
        self.conn.execute("DELETE FROM summaries WHERE paper_id=?", (pid,))
        self.conn.commit()

    def clear_papers(self):
        self.conn.execute("DELETE FROM papers")
        self.conn.execute("DELETE FROM summaries")
        self.conn.commit()

    def stats(self) -> dict:
        papers   = self.conn.execute("SELECT COUNT(*) FROM papers").fetchone()[0]
        summaries = self.conn.execute("SELECT COUNT(*) FROM summaries").fetchone()[0]
        searches  = self.conn.execute("SELECT COUNT(*) FROM searches").fetchone()[0]
        return {"total_papers": papers, "total_summaries": summaries,
                "total_searches": searches}


db = DB()


# ══════════════════════════════════════════════════════════════════════════════
# PAPER FINDER
# ══════════════════════════════════════════════════════════════════════════════

def search_arxiv(query: str, n: int) -> List[dict]:
    try:
        import arxiv
        search = arxiv.Search(
            query=query, max_results=n * 2,
            sort_by=arxiv.SortCriterion.Relevance
        )
        papers = []
        for r in search.results():
            papers.append({
                "id": r.entry_id.split("/")[-1],
                "title": r.title.strip(),
                "authors": [a.name for a in r.authors],
                "abstract": r.summary.strip(),
                "year": r.published.year,
                "url": r.pdf_url or "",
                "source": "arXiv",
                "citations": 0,
            })
        return papers
    except Exception as e:
        print(f"arXiv error: {e}")
        return []


def search_semantic_scholar(query: str, n: int) -> List[dict]:
    try:
        import requests
        api_key = os.getenv("SEMANTIC_SCHOLAR_API_KEY", "")
        params  = {"query": query, "limit": n * 2,
                   "fields": "title,authors,abstract,year,citationCount,externalIds"}
        headers = {"x-api-key": api_key} if api_key else {}

        resp = requests.get(
            "https://api.semanticscholar.org/graph/v1/paper/search",
            params=params, timeout=12, headers=headers
        )
        # If key is invalid/expired, retry without it
        if resp.status_code == 403 and api_key:
            resp = requests.get(
                "https://api.semanticscholar.org/graph/v1/paper/search",
                params=params, timeout=12
            )
        if resp.status_code != 200:
            return []
        papers = []
        for item in resp.json().get("data", []):
            pid = item.get("paperId", "unknown")
            ext = item.get("externalIds") or {}
            if ext.get("ArXiv"):
                pid = ext["ArXiv"]
            papers.append({
                "id": pid,
                "title": (item.get("title") or "Unknown").strip(),
                "authors": [a.get("name", "") for a in item.get("authors", [])],
                "abstract": (item.get("abstract") or "No abstract available.").strip(),
                "year": item.get("year") or 0,
                "citations": item.get("citationCount") or 0,
                "url": f"https://semanticscholar.org/paper/{pid}",
                "source": "Semantic Scholar",
            })
        return papers
    except Exception as e:
        print(f"Semantic Scholar error: {e}")
        return []


def find_papers(query: str, n: int) -> List[dict]:
    # Give each source a share so SS results aren't crowded out by arXiv
    half = max(n // 2, 5)
    all_papers = search_arxiv(query, half) + search_semantic_scholar(query, half)
    seen, unique = set(), []
    for p in all_papers:
        key = re.sub(r"\W+", "", p["title"].lower())[:60]
        if key not in seen and len(key) > 5:
            seen.add(key)
            unique.append(p)
    return unique[:n]


# ══════════════════════════════════════════════════════════════════════════════
# AI ENGINE  (Groq + Llama3)
# ══════════════════════════════════════════════════════════════════════════════

_groq_client = None

def _get_groq():
    global _groq_client
    if _groq_client is None:
        key = os.getenv("GROQ_API_KEY", "")
        if key:
            try:
                from groq import Groq
                _groq_client = Groq(api_key=key)
                print("✓ Groq AI engine initialized (Llama3)")
            except ImportError:
                print("! groq package not installed — run: pip install groq")
            except Exception as e:
                print(f"! Groq init error: {e}")
    return _groq_client

def ai_available() -> bool:
    return _get_groq() is not None

def llm_call(prompt: str,
             system: str = "You are a helpful expert research assistant.",
             model: str = "llama-3.1-8b-instant",
             max_tokens: int = 1024) -> Optional[str]:
    """Call Groq LLM. Returns None if unavailable or on error."""
    client = _get_groq()
    if not client:
        return None
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": prompt},
            ],
            temperature=0.35,
            max_tokens=max_tokens,
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        print(f"LLM error: {e}")
        return None


# ══════════════════════════════════════════════════════════════════════════════
# SUMMARIZER
# ══════════════════════════════════════════════════════════════════════════════

def _extractive_summarize(abstract: str) -> dict:
    """Fast extractive fallback when no AI key."""
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", abstract) if s.strip()]
    keywords = ["propose", "present", "introduce", "novel", "new", "show",
                "demonstrate", "achieve", "we ", "our ", "this paper"]
    scored = []
    for s in sentences:
        sl = s.lower()
        score = sum(2 for kw in keywords if kw in sl) + min(len(s.split()) / 5, 3)
        scored.append((score, s))
    scored.sort(reverse=True)
    tldr  = scored[0][1] if scored else sentences[0]
    words = tldr.split()
    tldr  = " ".join(words[:60]) + ("..." if len(words) > 60 else "")
    short = " ".join(sentences[:3])
    markers = ["we propose", "we present", "we introduce", "we show",
               "our approach", "our method", "we achieve", "results show",
               "outperform", "novel", "state-of-the-art"]
    keypoints = [s for s in sentences if any(m in s.lower() for m in markers)][:5]
    if not keypoints:
        step = max(1, len(sentences) // 5)
        keypoints = sentences[::step][:5]
    return {"tldr": tldr, "short": short, "full": abstract,
            "keypoints": keypoints, "ai_powered": False}


def summarize_paper(paper: dict) -> dict:
    abstract = paper.get("abstract", "").strip()
    title    = paper.get("title", "").strip()
    if not abstract or len(abstract) < 30:
        return {"tldr": "No abstract available.", "short": "No abstract available.",
                "full": "No abstract available.", "keypoints": [], "ai_powered": False}

    # ── Try Groq AI first ────────────────────────────────────────────────────
    if ai_available():
        prompt = f"""Analyze this research paper and respond in EXACTLY this format:

TLDR: <one powerful sentence, max 35 words, capturing the core contribution>
CONTRIBUTION_1: <first key contribution, one sentence>
CONTRIBUTION_2: <second key contribution, one sentence>
CONTRIBUTION_3: <third key contribution, one sentence>
FULL_SUMMARY: <3-4 sentences synthesizing the problem, method, results, and significance>

Paper Title: {title}
Abstract: {abstract[:1500]}"""

        result = llm_call(
            prompt,
            system="You are an expert research analyst. Extract precise, insightful information. Follow the format exactly.",
            model="llama-3.1-8b-instant",
            max_tokens=512,
        )
        if result:
            tldr = ""; keypoints = []; full = ""
            for line in result.splitlines():
                line = line.strip()
                if line.startswith("TLDR:"):
                    tldr = line[5:].strip()
                elif line.startswith("CONTRIBUTION_"):
                    kp = line.split(":", 1)[-1].strip()
                    if kp:
                        keypoints.append(kp)
                elif line.startswith("FULL_SUMMARY:"):
                    full = line[13:].strip()
                elif full and line and not line.startswith(("TLDR", "CONTRIB", "FULL")):
                    full += " " + line
            if tldr and keypoints:
                return {
                    "tldr":      tldr,
                    "short":     tldr + " " + (keypoints[0] if keypoints else ""),
                    "full":      full or abstract,
                    "keypoints": keypoints,
                    "ai_powered": True,
                }

    # ── Extractive fallback ──────────────────────────────────────────────────
    return _extractive_summarize(abstract)


# ══════════════════════════════════════════════════════════════════════════════
# PRESENTATION GENERATOR  — 14-slide academic layout
# ══════════════════════════════════════════════════════════════════════════════

_STOP_WORDS = {
    "the","a","an","of","in","for","and","or","to","is","are","with","that",
    "this","we","our","on","at","by","from","as","be","it","its","also","but",
    "not","have","has","was","were","using","used","based","which","paper",
    "show","can","when","their","they","been","more","than","these","such",
    "both","each","after","before","other","into","through","during","over",
}

def _extract_pptx_content(paper: dict, summary: dict) -> dict:
    """Parse abstract + summary into structured slide content."""
    abstract = paper.get("abstract", "")
    title    = paper.get("title", "")
    authors  = paper.get("authors", [])

    sents = [s.strip() for s in re.split(r"(?<=[.!?])\s+", abstract) if len(s.strip()) > 20]

    def pick(keywords, n=3):
        hits = [s for s in sents if any(k in s.lower() for k in keywords)]
        return hits[:n] if hits else sents[:min(n, len(sents))]

    # Extract top technical keywords as module names
    text  = (title + " " + abstract).lower()
    words = re.findall(r"\b[a-z][a-z]{4,}\b", text)
    freq  = {}
    for w in words:
        if w not in _STOP_WORDS:
            freq[w] = freq.get(w, 0) + 1
    modules = sorted(freq, key=freq.get, reverse=True)[:8]

    auth_str = ", ".join(authors[:3])
    if len(authors) > 3:
        auth_str += " et al."

    return {
        "title":      title,
        "authors":    auth_str,
        "year":       paper.get("year", ""),
        "source":     paper.get("source", ""),
        "url":        paper.get("url", ""),
        "tldr":       summary.get("tldr", ""),
        "keypoints":  summary.get("keypoints", []),
        "abstract":   abstract,
        "intro":      sents[:3],
        "problems":   pick(["however","challenge","limitation","problem","difficult",
                            "lack","existing","current","traditional","previous"], 3),
        "objectives": pick(["aim","objective","goal","propose","present","introduce",
                            "develop","design","investigate","focus"], 3),
        "method":     pick(["propose","present","method","approach","framework",
                            "algorithm","model","architecture","system","technique"], 4),
        "results":    pick(["achieve","result","accuracy","performance","improve",
                            "outperform","show","demonstrate","experiment",
                            "state-of-the-art","evaluate"], 4),
        "future":     pick(["future","extend","plan","further","next","potential",
                            "could","will","limitation","open"], 3),
        "modules":    modules,
    }


def generate_pptx(paper: dict, summary: dict) -> Optional[str]:
    try:
        from pptx import Presentation as PPT
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        prs = PPT()
        c   = _extract_pptx_content(paper, summary)

        # ── colour palette ──────────────────────────────────────────────────
        NAVY   = RGBColor(0x1e, 0x1b, 0x4b)
        PURPLE = RGBColor(0x7c, 0x3a, 0xed)
        WHITE  = RGBColor(0xff, 0xff, 0xff)

        def _style_title(sl):
            """Make the title placeholder bold + navy."""
            try:
                tf = sl.shapes.title.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.bold  = True
                        run.font.color.rgb = NAVY
                    para.font.bold = True
            except Exception:
                pass

        def add_slide(title_text: str, bullets: list):
            """Add a Title+Content slide and populate with bullet lines."""
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title_text
            _style_title(sl)

            tf = sl.placeholders[1].text_frame
            tf.clear()
            first = True
            for line in bullets:
                line = str(line)[:300]
                if not line:
                    continue
                if first:
                    p = tf.paragraphs[0]; first = False
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                # indent sub-items
                if line.startswith("  ") or line.startswith("\t"):
                    p.level = 1
            return sl

        m = c["modules"]

        # ── 1. Title ─────────────────────────────────────────────────────────
        sl1 = prs.slides.add_slide(prs.slide_layouts[0])
        sl1.shapes.title.text = c["title"]
        sl1.placeholders[1].text = (
            f"{c['authors']}\n"
            f"{c['year']}  ·  {c['source']}\n\n"
            f"Submitted for Academic Presentation"
        )
        _style_title(sl1)

        # ── 2. Agenda ─────────────────────────────────────────────────────────
        add_slide("Agenda", [
            "1.  Introduction",
            "2.  Literature Survey",
            "3.  Problem Statement / Limitations",
            "4.  Objective",
            "5.  Proposed System",
            "6.  Architecture Diagram",
            "7.  Modules",
            "8.  System Requirements",
            "9.  Results",
            "10. Conclusion",
            "11. Future Enhancements",
            "12. References",
        ])

        # ── 3. Introduction ───────────────────────────────────────────────────
        intro_lines = [f"• {s}" for s in c["intro"]]
        if not intro_lines:
            intro_lines = [f"• {c['abstract'][:300]}"]
        add_slide("Introduction", intro_lines)

        # ── 4. Literature Survey ──────────────────────────────────────────────
        lit = [
            f"• Existing research in {m[0] if m else 'this domain'} provides foundational methods",
            f"• Prior work on {m[1] if len(m)>1 else 'related topics'} highlights key challenges",
        ]
        if c["problems"]:
            lit.append(f"• {c['problems'][0]}")
        lit += [
            f"• Gap identified in current literature on {m[2] if len(m)>2 else 'the proposed approach'}",
            f"• This study extends prior work on {', '.join(m[:3]) if m else 'the topic'}",
        ]
        add_slide("Literature Survey", lit)

        # ── 5. Problem Statement / Limitations ───────────────────────────────
        prob = [f"• {s}" for s in c["problems"]]
        if not prob:
            prob = [
                f"• Existing {m[0] if m else 'approaches'} exhibit performance limitations",
                "• Current solutions do not generalize well to real-world scenarios",
                "• Lack of efficiency and scalability in prior methods",
            ]
        add_slide("Problem Statement / Limitations", prob)

        # ── 6. Objective ──────────────────────────────────────────────────────
        obj = []
        if c["tldr"]:
            obj.append(f"• {c['tldr']}")
        obj += [f"• {s}" for s in c["objectives"][:3]]
        if not obj:
            obj = [
                f"• To design a novel solution for {c['title'][:70]}",
                "• To evaluate performance against state-of-the-art baselines",
                "• To demonstrate practical applicability on real datasets",
            ]
        add_slide("Objective", obj)

        # ── 7. Proposed System ────────────────────────────────────────────────
        meth = [f"• {s}" for s in c["method"][:4]]
        if not meth:
            meth = [
                f"• Proposed: a novel {m[0] if m else 'system'}-based approach",
                "• End-to-end trainable framework",
                "• Addresses identified limitations through innovative design",
            ]
        add_slide("Proposed System", meth)

        # ── 8. Architecture Diagram ───────────────────────────────────────────
        arch = [
            "High-Level System Architecture:",
            "",
            f"  [ Input Data / Dataset ]",
            f"         ↓",
            f"  [ {m[0].title() if m else 'Pre-processing'} Module ]",
            f"         ↓",
            f"  [ {m[1].title() if len(m)>1 else 'Core'} Engine  ←→  {m[2].title() if len(m)>2 else 'Analysis'} Layer ]",
            f"         ↓",
            f"  [ {m[3].title() if len(m)>3 else 'Output'} / Results ]",
            "",
            "  * See the paper for the complete architecture diagram.",
        ]
        add_slide("Architecture Diagram", arch)

        # ── 9. Modules ────────────────────────────────────────────────────────
        mods = [f"• {w.replace('-', ' ').title()} Module" for w in m[:7]]
        if not mods:
            mods = ["• Input Processing Module", "• Core Analysis Module", "• Output Generation Module"]
        add_slide("Modules", mods)

        # ── 10. System Requirements ───────────────────────────────────────────
        add_slide("System Requirements", [
            "Hardware Requirements:",
            "  • CPU: Intel Core i5 / AMD Ryzen 5 or higher",
            "  • RAM: 8 GB minimum  (16 GB recommended)",
            "  • Storage: 10 GB free disk space",
            "  • GPU: NVIDIA GPU with CUDA support (for DL models)",
            "",
            "Software Requirements:",
            "  • OS: Windows 10 / Ubuntu 20.04 / macOS 12+",
            "  • Python 3.8+  |  pip package manager",
            f"  • Key libraries: {', '.join(m[:5]) if m else 'standard ML/DL stack'}",
        ])

        # ── 11. Results ───────────────────────────────────────────────────────
        res = [f"• {s}" for s in c["results"][:4]]
        if not res:
            res = [
                "• Proposed method evaluated on standard benchmark datasets",
                "• Achieves state-of-the-art performance on key evaluation metrics",
                "• Significant improvement in accuracy/efficiency over baselines",
            ]
        add_slide("Results", res)

        # ── 12. Conclusion ────────────────────────────────────────────────────
        conc = []
        if c["tldr"]:
            conc.append(f"• {c['tldr']}")
        conc += [f"• {kp}" for kp in c["keypoints"][:3]]
        if not conc:
            conc = [
                "• Successfully addressed the identified research challenges",
                "• Proposed approach demonstrates superior performance",
                "• Work contributes meaningfully to the research community",
            ]
        add_slide("Conclusion", conc)

        # ── 13. Future Enhancements ───────────────────────────────────────────
        fut = [f"• {s}" for s in c["future"][:3]]
        if not fut:
            fut = [
                "• Extend to larger and more diverse real-world datasets",
                "• Optimize for real-time deployment and edge computing",
                "• Investigate cross-domain transfer learning capabilities",
                "• Explore integration with complementary methods",
            ]
        add_slide("Future Enhancements", fut)

        # ── 14. References ────────────────────────────────────────────────────
        add_slide("References", [
            f"[1] {c['authors']} ({c['year']}).",
            f"    {c['title'][:100]}.",
            f"    {c['source']}.",
            f"    URL: {c['url']}",
            "",
            "For the complete reference list, please refer to the original paper.",
        ])

        # ── Save ──────────────────────────────────────────────────────────────
        safe = re.sub(r"[^\w\s-]", "", paper["title"])[:40].strip().replace(" ", "_")
        path = str(PPTX_DIR / f"{safe}.pptx")
        prs.save(path)
        return path

    except Exception as e:
        print(f"PPTX error: {e}")
        import traceback; traceback.print_exc()
        return None


# ══════════════════════════════════════════════════════════════════════════════
# KNOWLEDGE GRAPH
# ══════════════════════════════════════════════════════════════════════════════

def build_knowledge_graph(papers: List[dict]) -> dict:
    stop = {"the", "a", "an", "of", "in", "for", "and", "or", "to", "is",
            "are", "with", "that", "this", "we", "our", "on", "at", "by",
            "from", "as", "be", "it", "its", "also", "but", "not", "have",
            "has", "was", "were", "using", "used", "based", "which"}

    nodes, edges = [], []
    concept_map: Dict[str, List[str]] = {}

    for paper in papers:
        text = (paper["title"] + " " + paper.get("abstract", "")).lower()
        words = re.findall(r"\b[a-z]{4,}\b", text)
        freq: Dict[str, int] = {}
        for w in words:
            if w not in stop:
                freq[w] = freq.get(w, 0) + 1
        top_concepts = sorted(freq, key=freq.get, reverse=True)[:8]

        nodes.append({
            "id": paper["id"],
            "label": paper["title"][:50],
            "type": "paper",
            "year": paper.get("year", 0),
            "concepts": top_concepts,
        })

        for c in top_concepts:
            concept_map.setdefault(c, []).append(paper["id"])

    # Edges between papers sharing concepts
    shared_concepts: Dict[str, List[str]] = {
        k: v for k, v in concept_map.items() if len(v) > 1
    }
    added_edges = set()
    for concept, pids in shared_concepts.items():
        for i in range(len(pids)):
            for j in range(i + 1, len(pids)):
                edge_key = tuple(sorted([pids[i], pids[j]]))
                if edge_key not in added_edges:
                    edges.append({
                        "source": pids[i],
                        "target": pids[j],
                        "label": concept,
                    })
                    added_edges.add(edge_key)

    return {
        "nodes": nodes,
        "edges": edges,
        "shared_concepts": [
            {"concept": k, "papers": v, "count": len(v)}
            for k, v in sorted(shared_concepts.items(),
                               key=lambda x: len(x[1]), reverse=True)[:20]
        ],
    }


# ══════════════════════════════════════════════════════════════════════════════
# GAP DISCOVERY  (lightweight, no ML needed)
# ══════════════════════════════════════════════════════════════════════════════

def discover_gaps(papers: List[dict]) -> List[dict]:
    gaps = []
    all_concepts: Dict[str, List[str]] = {}

    stop = {"the","a","an","of","in","for","and","or","to","is","are","with",
            "that","this","we","our","on","at","by","from","as","be","it",
            "its","also","but","not","have","has","was","were","using","used"}

    for paper in papers:
        text = (paper["title"] + " " + paper.get("abstract","")).lower()
        words = re.findall(r"\b[a-z]{4,}\b", text)
        for w in set(words):
            if w not in stop:
                all_concepts.setdefault(w, []).append(paper["id"])

    # Gap 1: concepts appearing in only one paper (underexplored)
    single_paper_concepts = [(c, pids) for c, pids in all_concepts.items()
                              if len(pids) == 1 and len(c) > 5]
    single_paper_concepts.sort(key=lambda x: len(x[0]), reverse=True)
    for concept, pids in single_paper_concepts[:5]:
        gaps.append({
            "gap_id": f"gap_underexplored_{concept}",
            "type": "Underexplored Area",
            "title": f"Underexplored: '{concept}'",
            "description": f"The concept '{concept}' appears in only 1 paper. High opportunity for deeper research.",
            "confidence": 0.72,
            "impact": "Medium",
            "related_papers": pids,
        })

    # Gap 2: temporal gap — old papers with no recent follow-up
    years = [(p.get("year", 0), p["id"], p["title"]) for p in papers if p.get("year")]
    if years:
        min_year = min(y[0] for y in years)
        max_year = max(y[0] for y in years)
        old_papers = [y for y in years if y[0] <= min_year + 2]
        if old_papers and max_year - min_year > 3:
            for yr, pid, title in old_papers[:3]:
                gaps.append({
                    "gap_id": f"gap_temporal_{pid}",
                    "type": "Temporal Gap",
                    "title": f"Old paper with potential for update ({yr})",
                    "description": f"'{title[:60]}' from {yr} may need re-evaluation with modern methods.",
                    "confidence": 0.68,
                    "impact": "High",
                    "related_papers": [pid],
                })

    # Gap 3: cross-domain — detect when papers from different areas share concepts
    sources = set(p.get("source","") for p in papers)
    for concept, pids in all_concepts.items():
        if len(pids) >= 2 and len(concept) > 6:
            paper_sources = [next((p.get("source","") for p in papers if p["id"] == pid), "") for pid in pids]
            if len(set(paper_sources)) > 1:
                gaps.append({
                    "gap_id": f"gap_crossdomain_{concept}",
                    "type": "Cross-Domain Opportunity",
                    "title": f"Cross-domain link via '{concept}'",
                    "description": f"Concept '{concept}' appears across {len(pids)} papers from different sources. Strong cross-domain opportunity.",
                    "confidence": 0.81,
                    "impact": "Breakthrough",
                    "related_papers": pids[:4],
                })
                if len(gaps) >= 12:
                    break

    # Gap 4: methodological gap — look for "baseline" without "comparison"
    for paper in papers:
        abstract = paper.get("abstract","").lower()
        if "baseline" in abstract and "comparison" not in abstract:
            gaps.append({
                "gap_id": f"gap_method_{paper['id']}",
                "type": "Methodological Gap",
                "title": f"Missing comparison study",
                "description": f"'{paper['title'][:60]}' mentions baselines but lacks systematic comparison.",
                "confidence": 0.75,
                "impact": "Medium",
                "related_papers": [paper["id"]],
            })

    return gaps[:15]


def discover_gaps_ai(papers: List[dict]) -> List[dict]:
    """AI-powered research gap discovery using Groq."""
    if not ai_available() or len(papers) < 2:
        return discover_gaps(papers)

    paper_lines = []
    for p in papers[:18]:
        paper_lines.append(f"[{p.get('year','')}] {p['title']}: {p.get('abstract','')[:220]}")
    papers_text = "\n".join(paper_lines)

    prompt = f"""You are a senior research strategist. Analyze these papers and identify 6 specific, actionable research gaps.

PAPERS:
{papers_text}

For each gap respond in EXACTLY this format (no deviation):
GAP_TYPE: <one of: Underexplored Area | Temporal Gap | Methodological Gap | Cross-Domain Opportunity | Theoretical Gap | Missing Dataset>
GAP_TITLE: <short specific title, max 12 words>
GAP_DESCRIPTION: <2-3 sentences: what's missing, why it matters, what opportunity it creates>
IMPACT: <High | Breakthrough | Medium>
CONFIDENCE: <0.65-0.95>
---
(repeat for all 6 gaps)"""

    result = llm_call(
        prompt,
        system="You are a world-class research strategist identifying high-value gaps in scientific literature. Be specific, novel, and actionable.",
        model="llama-3.1-8b-instant",
        max_tokens=1600,
    )
    if not result:
        return discover_gaps(papers)

    gaps = []
    current: Dict[str, Any] = {}
    for line in result.splitlines():
        line = line.strip()
        if line.startswith("GAP_TYPE:"):
            if current.get("title"):
                gaps.append(current); current = {}
            current["type"] = line[9:].strip()
        elif line.startswith("GAP_TITLE:"):
            current["title"] = line[10:].strip()
        elif line.startswith("GAP_DESCRIPTION:"):
            current["description"] = line[16:].strip()
        elif line.startswith("IMPACT:"):
            current["impact"] = line[7:].strip()
        elif line.startswith("CONFIDENCE:"):
            try:    current["confidence"] = float(line[11:].strip())
            except: current["confidence"] = 0.78
        elif line == "---" and current.get("title"):
            gaps.append(current); current = {}
        elif current.get("description") and line and not any(
            line.startswith(k) for k in ["GAP_","IMPACT","CONFIDENCE","---"]
        ):
            current["description"] += " " + line

    if current.get("title"):
        gaps.append(current)

    formatted = []
    for i, g in enumerate(gaps[:8]):
        formatted.append({
            "gap_id":        f"ai_gap_{i}",
            "type":          g.get("type", "Research Gap"),
            "title":         g.get("title", "Unknown gap"),
            "description":   g.get("description", ""),
            "confidence":    g.get("confidence", 0.78),
            "impact":        g.get("impact", "High"),
            "related_papers": [papers[min(i, len(papers)-1)]["id"]],
            "ai_powered":    True,
        })
    return formatted if formatted else discover_gaps(papers)


# ══════════════════════════════════════════════════════════════════════════════
# SEMANTIC ENGINE  (sentence-transformers)
# ══════════════════════════════════════════════════════════════════════════════

_embedder = None

def _get_embedder():
    global _embedder
    if _embedder is None:
        try:
            from sentence_transformers import SentenceTransformer
            _embedder = SentenceTransformer("all-MiniLM-L6-v2")
            print("✓ Semantic embedder loaded (all-MiniLM-L6-v2)")
        except Exception as e:
            print(f"! Embedder load error: {e}")
    return _embedder

def semantic_available() -> bool:
    return _get_embedder() is not None

def embed_texts(texts: List[str]):
    model = _get_embedder()
    if not model:
        return None
    try:
        return model.encode(texts, convert_to_numpy=True, show_progress_bar=False)
    except Exception as e:
        print(f"Embed error: {e}")
        return None

def cosine_sim(a, b) -> float:
    import numpy as np
    a, b = np.array(a), np.array(b)
    denom = (np.linalg.norm(a) * np.linalg.norm(b))
    return float(np.dot(a, b) / denom) if denom > 0 else 0.0

def semantic_search(query: str, papers: List[dict], top_k: int = 10) -> List[dict]:
    """Rank papers by semantic similarity to query."""
    if not papers:
        return []
    model = _get_embedder()
    if not model:
        return papers[:top_k]
    texts = [f"{p['title']} {p.get('abstract','')[:300]}" for p in papers]
    try:
        vecs   = embed_texts([query] + texts)
        q_vec  = vecs[0]
        p_vecs = vecs[1:]
        scored = [(cosine_sim(q_vec, pv), p) for pv, p in zip(p_vecs, papers)]
        scored.sort(key=lambda x: x[0], reverse=True)
        return [dict(**p, similarity=round(s, 3)) for s, p in scored[:top_k]]
    except Exception as e:
        print(f"Semantic search error: {e}")
        return papers[:top_k]

def find_similar_papers(paper: dict, all_papers: List[dict], top_k: int = 5) -> List[dict]:
    """Find papers most similar to a given paper using embeddings."""
    if not all_papers:
        return []
    others = [p for p in all_papers if p["id"] != paper["id"]]
    if not others:
        return []
    query = f"{paper['title']} {paper.get('abstract','')[:300]}"
    return semantic_search(query, others, top_k)


# ══════════════════════════════════════════════════════════════════════════════
# TREND ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════

def analyze_trends(papers: List[dict]) -> dict:
    """Analyze publication trends and topic evolution over time."""
    by_year: Dict[int, List[dict]] = {}
    for p in papers:
        y = p.get("year", 0)
        if y and y > 1990:
            by_year.setdefault(y, []).append(p)

    if not by_year:
        return {"years": [], "counts": [], "top_concepts_by_era": {}, "velocity": {}}

    years_sorted = sorted(by_year.keys())

    # Top concepts per era (split into early/mid/recent thirds)
    n = len(years_sorted)
    eras = {
        "Early":  years_sorted[:max(1, n//3)],
        "Middle": years_sorted[max(1, n//3):max(2, 2*n//3)],
        "Recent": years_sorted[max(2, 2*n//3):],
    }
    stop = {"the","a","an","of","in","and","or","to","is","are","with","that",
            "this","we","our","on","at","by","from","as","be","it","its","also",
            "but","not","have","has","was","were","using","used","paper","show",
            "propose","present","approach","method","results","study","based",
            "these","their","which","such","more","when","been","each","other"}

    era_concepts: Dict[str, List[str]] = {}
    for era_name, era_years in eras.items():
        era_papers = [p for y in era_years for p in by_year.get(y, [])]
        if not era_papers:
            continue
        freq: Dict[str, int] = {}
        for p in era_papers:
            txt = (p["title"] + " " + p.get("abstract","")[:200]).lower()
            for w in re.findall(r"\b[a-z]{5,}\b", txt):
                if w not in stop:
                    freq[w] = freq.get(w, 0) + 1
        top = sorted(freq, key=freq.get, reverse=True)[:8]  # type: ignore[arg-type]
        era_concepts[era_name] = top

    # Publication velocity (papers per year)
    velocity = {str(y): len(ps) for y, ps in by_year.items()}

    # Rising topics: concepts that appear more in recent years than early years
    rising: List[Dict[str, Any]] = []
    early_papers  = [p for y in eras["Early"]  for p in by_year.get(y, [])]
    recent_papers = [p for y in eras["Recent"] for p in by_year.get(y, [])]

    def concept_freq(paper_list: List[dict]) -> Dict[str, int]:
        f: Dict[str, int] = {}
        for p in paper_list:
            txt = (p["title"] + " " + p.get("abstract","")[:200]).lower()
            for w in re.findall(r"\b[a-z]{5,}\b", txt):
                if w not in stop:
                    f[w] = f.get(w, 0) + 1
        return f

    ef = concept_freq(early_papers)
    rf = concept_freq(recent_papers)
    if ef and rf:
        for concept, rcnt in rf.items():
            ecnt = ef.get(concept, 0)
            if rcnt > 1 and rcnt > ecnt * 1.5:
                rising.append({"concept": concept, "early": ecnt, "recent": rcnt,
                                "growth": round(rcnt / max(ecnt, 1), 1)})
        rising.sort(key=lambda x: x["growth"], reverse=True)

    return {
        "years":               years_sorted,
        "counts":              [len(by_year[y]) for y in years_sorted],
        "top_concepts_by_era": era_concepts,
        "velocity":            velocity,
        "rising_topics":       rising[:8],
        "total_papers":        len(papers),
        "year_range":          [min(years_sorted), max(years_sorted)],
    }


# ══════════════════════════════════════════════════════════════════════════════
# BIBTEX EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def paper_to_bibtex(paper: dict) -> str:
    """Convert a paper dict to BibTeX entry."""
    pid    = re.sub(r"\W+", "", paper["id"])[:20]
    auth   = paper.get("authors", [])
    author_str = " and ".join(auth[:6]) if auth else "Unknown"
    title  = paper.get("title", "Unknown").replace("{", "\\{").replace("}", "\\}")
    year   = paper.get("year", "")
    url    = paper.get("url", "")
    src    = paper.get("source", "")
    journal = "arXiv preprint" if "arXiv" in src else "Semantic Scholar"
    entry  = (
        f"@article{{{pid},\n"
        f"  author  = {{{author_str}}},\n"
        f"  title   = {{{title}}},\n"
        f"  year    = {{{year}}},\n"
        f"  journal = {{{journal}}},\n"
        f"  url     = {{{url}}},\n"
        f"}}"
    )
    return entry


# ══════════════════════════════════════════════════════════════════════════════
# REQUEST / RESPONSE MODELS
# ══════════════════════════════════════════════════════════════════════════════

class SearchRequest(BaseModel):
    query: str
    max_results: int = 10

class SummarizeRequest(BaseModel):
    paper_id: str

class PresentRequest(BaseModel):
    paper_id: str

class GapRequest(BaseModel):
    paper_ids: Optional[List[str]] = None   # None = use all saved papers

class KnowledgeGraphRequest(BaseModel):
    paper_ids: Optional[List[str]] = None

class ChatRequest(BaseModel):
    message: str
    paper_ids: Optional[List[str]] = None
    history:   Optional[List[Dict[str, str]]] = None  # [{role, content}, ...]

class ReviewRequest(BaseModel):
    topic:     str
    paper_ids: Optional[List[str]] = None

class IdeasRequest(BaseModel):
    topic:     str = ""
    paper_ids: Optional[List[str]] = None

class SemanticSearchRequest(BaseModel):
    query:     str
    top_k:     int = 10


# ══════════════════════════════════════════════════════════════════════════════
# API ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/")
def root():
    return {"message": "ScholarGenie API", "version": "1.0.0",
            "docs": "/docs", "status": "running"}


@app.get("/health")
def health():
    s = db.stats()
    return {"status": "healthy", "database": "sqlite", **s}


# ── Papers ────────────────────────────────────────────────────────────────────

@app.post("/api/search")
def search(req: SearchRequest):
    """Search papers from arXiv and Semantic Scholar."""
    if not req.query.strip():
        raise HTTPException(400, "Query cannot be empty")

    papers = find_papers(req.query.strip(), max(1, min(req.max_results, 50)))

    for p in papers:
        db.save_paper(p)
    db.log_search(req.query, len(papers))

    return {"query": req.query, "count": len(papers), "papers": papers}


@app.get("/api/papers")
def get_papers():
    """Get all saved papers."""
    papers = db.get_papers()
    return {"papers": papers, "count": len(papers)}


@app.get("/api/papers/{paper_id}")
def get_paper(paper_id: str):
    """Get a single paper by ID."""
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found")
    summary = db.get_summary(paper_id)
    return {"paper": paper, "summary": summary}


# ── Summarize ─────────────────────────────────────────────────────────────────

@app.post("/api/summarize")
def summarize(req: SummarizeRequest):
    """Summarize a paper."""
    paper = db.get_paper(req.paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found. Search for it first.")

    summary = summarize_paper(paper)
    db.save_summary(req.paper_id, summary)

    return {"paper_id": req.paper_id, "title": paper["title"], "summary": summary}


@app.post("/api/summarize/batch")
def summarize_batch(paper_ids: List[str]):
    """Summarize multiple papers."""
    results = []
    for pid in paper_ids:
        paper = db.get_paper(pid)
        if paper:
            summary = summarize_paper(paper)
            db.save_summary(pid, summary)
            results.append({"paper_id": pid, "title": paper["title"], "summary": summary})
    return {"count": len(results), "summaries": results}


# ── Presentation ──────────────────────────────────────────────────────────────

@app.post("/api/present")
def present(req: PresentRequest):
    """Generate PowerPoint for a paper."""
    paper = db.get_paper(req.paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found.")

    summary = db.get_summary(req.paper_id) or summarize_paper(paper)
    path    = generate_pptx(paper, summary)

    if not path:
        raise HTTPException(500, "Failed to generate presentation.")

    return {
        "paper_id":  req.paper_id,
        "title":     paper["title"],
        "file_path": path,
        "download":  f"/api/present/download/{req.paper_id}",
        "slides":    14,
    }


@app.get("/api/present/download/{paper_id}")
def download_pptx(paper_id: str):
    """Download generated PowerPoint."""
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found.")
    safe  = re.sub(r"[^\w\s-]", "", paper["title"])[:40].strip().replace(" ", "_")
    path  = PPTX_DIR / f"{safe}.pptx"
    if not path.exists():
        raise HTTPException(404, "Presentation not generated yet. Call /api/present first.")
    return FileResponse(str(path), media_type="application/vnd.ms-powerpoint",
                        filename=f"{safe}.pptx")


# ── Knowledge Graph ───────────────────────────────────────────────────────────

@app.post("/api/knowledge-graph")
def knowledge_graph(req: KnowledgeGraphRequest = None):
    """Build knowledge graph from papers. Uses Neo4j if configured, else in-memory."""
    if req is None:
        req = KnowledgeGraphRequest()
    if req.paper_ids:
        papers = [db.get_paper(pid) for pid in req.paper_ids]
        papers = [p for p in papers if p]
    else:
        papers = db.get_papers()

    if not papers:
        raise HTTPException(404, "No papers found.")

    # Try Neo4j first; fall back to in-memory builder
    try:
        from graph import get_neo4j
        neo4j = get_neo4j()
        if neo4j.available:
            graph = neo4j.build_graph(papers)
            return {
                "paper_count": len(papers),
                "graph": graph,
                "backend": "neo4j",
            }
    except Exception as e:
        print(f"! Neo4j graph build failed, falling back: {e}")

    graph = build_knowledge_graph(papers)
    return {
        "paper_count": len(papers),
        "graph": graph,
        "backend": "memory",
    }


@app.get("/api/neo4j-status")
def neo4j_status():
    """Check if Neo4j is configured and connected."""
    try:
        from graph import get_neo4j
        neo4j = get_neo4j()
        return {
            "available": neo4j.available,
            "node_count": neo4j.node_count() if neo4j.available else 0,
        }
    except Exception:
        return {"available": False, "node_count": 0}


# ── Gap Discovery ─────────────────────────────────────────────────────────────

@app.post("/api/gaps")
def gap_discovery(req: GapRequest):
    """Discover research gaps (AI-powered when GROQ_API_KEY is set)."""
    if req.paper_ids:
        papers = [db.get_paper(pid) for pid in req.paper_ids]
        papers = [p for p in papers if p]
    else:
        papers = db.get_papers()

    if not papers:
        raise HTTPException(404, "No papers found. Search first.")

    gaps = discover_gaps_ai(papers) if ai_available() else discover_gaps(papers)
    return {
        "paper_count": len(papers),
        "gap_count":   len(gaps),
        "gaps":        gaps,
        "ai_powered":  ai_available(),
    }


# ── AI Chat ───────────────────────────────────────────────────────────────────

@app.post("/api/chat")
def chat_with_papers(req: ChatRequest):
    """Chat with your paper library using AI."""
    if not ai_available():
        raise HTTPException(400, "AI unavailable. Add GROQ_API_KEY to your .env file.")

    if req.paper_ids:
        papers = [db.get_paper(pid) for pid in req.paper_ids]
        papers = [p for p in papers if p]
    else:
        papers = db.get_papers()[:20]

    if not papers:
        raise HTTPException(404, "Library is empty. Search for papers first.")

    paper_ctx = []
    for p in papers[:15]:
        paper_ctx.append(
            f"Title: {p['title']} ({p.get('year','')})\n"
            f"Authors: {', '.join(p.get('authors',[])[:3])}\n"
            f"Abstract: {p.get('abstract','')[:350]}"
        )

    system = (
        f"You are ScholarGenie, an expert AI research assistant. "
        f"You have access to the user's library of {len(papers)} research papers.\n\n"
        f"LIBRARY:\n{'---'.join(paper_ctx)}\n\n"
        f"Answer questions based on these papers. Cite by [Title, Year]. "
        f"Be specific and insightful. If a topic isn't in the library, suggest searching for it."
    )

    messages: List[Dict[str, str]] = []
    if req.history:
        messages.extend(req.history[-8:])
    messages.append({"role": "user", "content": req.message})

    try:
        from groq import Groq
        client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        resp   = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "system", "content": system}] + messages,
            temperature=0.4,
            max_tokens=900,
        )
        reply = resp.choices[0].message.content.strip()
    except Exception as e:
        raise HTTPException(500, f"Chat error: {e}")

    return {
        "message":     req.message,
        "reply":       reply,
        "papers_used": len(papers),
        "ai_model":    "llama-3.1-8b-instant (Groq)",
    }


# ── Literature Review ─────────────────────────────────────────────────────────

@app.post("/api/review")
def literature_review(req: ReviewRequest):
    """Generate a full AI-powered literature review."""
    if not ai_available():
        raise HTTPException(400, "AI unavailable. Add GROQ_API_KEY to your .env file.")

    if req.paper_ids:
        papers = [db.get_paper(pid) for pid in req.paper_ids]
        papers = [p for p in papers if p]
    else:
        papers = db.get_papers()

    if not papers:
        raise HTTPException(404, "No papers found. Search for papers first.")

    paper_ctx = []
    for p in papers[:20]:
        authors_str = ", ".join(p.get("authors", [])[:2])
        paper_ctx.append(
            f"[{authors_str}, {p.get('year','')}] \"{p['title']}\"\n"
            f"  {p.get('abstract','')[:300]}"
        )
    papers_text = "\n\n".join(paper_ctx)

    prompt = f"""Write a comprehensive academic literature review on: "{req.topic}"

Using these {len(papers)} papers as sources:
{papers_text}

Structure the review with these sections. Use academic language. Cite as [Author, Year].

## 1. Introduction
(Scope, significance of the topic, 2-3 paragraphs)

## 2. Background & Foundational Work
(Historical context and seminal contributions, 2-3 paragraphs)

## 3. Current Approaches & Methodologies
(Main research directions, techniques, frameworks, 3-4 paragraphs)

## 4. Critical Analysis
(Strengths, limitations, and contradictions in current literature, 2-3 paragraphs)

## 5. Identified Research Gaps
(What is still missing or understudied, 2 paragraphs)

## 6. Conclusion & Future Directions
(Synthesis and promising directions, 1-2 paragraphs)"""

    result = llm_call(
        prompt,
        system="You are a senior academic researcher writing a thorough, well-structured literature review. Cite all papers appropriately.",
        model="llama-3.1-8b-instant",
        max_tokens=2500,
    )
    if not result:
        raise HTTPException(500, "Failed to generate literature review.")

    return {
        "topic":       req.topic,
        "paper_count": len(papers),
        "review":      result,
        "word_count":  len(result.split()),
        "ai_model":    "llama-3.3-70b-versatile (Groq)",
    }


# ── Research Ideas ────────────────────────────────────────────────────────────

@app.post("/api/ideas")
def research_ideas(req: IdeasRequest):
    """Generate novel AI-powered research ideas and hypotheses."""
    if not ai_available():
        raise HTTPException(400, "AI unavailable. Add GROQ_API_KEY to your .env file.")

    if req.paper_ids:
        papers = [db.get_paper(pid) for pid in req.paper_ids]
        papers = [p for p in papers if p]
    else:
        papers = db.get_papers()

    if not papers:
        raise HTTPException(404, "No papers found. Search first.")

    paper_lines = []
    for p in papers[:15]:
        paper_lines.append(f"[{p.get('year','')}] {p['title']}: {p.get('abstract','')[:200]}")
    papers_text = "\n".join(paper_lines)
    topic_hint  = f" Focus on: {req.topic}." if req.topic else ""

    prompt = f"""Based on these research papers, generate 5 novel and feasible research ideas.{topic_hint}

PAPERS:
{papers_text}

For each idea, use EXACTLY this format:
IDEA_TITLE: <catchy, specific title (max 12 words)>
HYPOTHESIS: <one testable scientific hypothesis sentence>
METHODOLOGY: <2-3 sentences on approach, tools, or experiments>
NOVELTY: <what makes this new vs existing work>
IMPACT: <High | Very High | Breakthrough>
DIFFICULTY: <Easy | Moderate | Hard | Very Hard>
---
(repeat for all 5 ideas)"""

    result = llm_call(
        prompt,
        system="You are a visionary research scientist. Generate creative, specific, and genuinely novel research ideas that build on the provided literature.",
        model="llama-3.1-8b-instant",
        max_tokens=1800,
    )
    if not result:
        raise HTTPException(500, "Failed to generate ideas.")

    ideas = []
    current: Dict[str, Any] = {}
    for line in result.splitlines():
        line = line.strip()
        if line.startswith("IDEA_TITLE:"):
            if current.get("title"): ideas.append(current); current = {}
            current["title"] = line[11:].strip()
        elif line.startswith("HYPOTHESIS:"):
            current["hypothesis"]   = line[11:].strip()
        elif line.startswith("METHODOLOGY:"):
            current["methodology"]  = line[12:].strip()
        elif line.startswith("NOVELTY:"):
            current["novelty"]      = line[8:].strip()
        elif line.startswith("IMPACT:"):
            current["impact"]       = line[7:].strip()
        elif line.startswith("DIFFICULTY:"):
            current["difficulty"]   = line[11:].strip()
        elif line == "---" and current.get("title"):
            ideas.append(current); current = {}
        elif current.get("title") and line and not any(
            line.startswith(k) for k in ["IDEA","HYPO","METHOD","NOVELTY","IMPACT","DIFF","---"]
        ):
            for field in ["methodology", "novelty", "hypothesis"]:
                if field in current and not current.get(field, "").endswith("."):
                    current[field] = (current.get(field, "") + " " + line).strip()
                    break
    if current.get("title"):
        ideas.append(current)

    return {
        "topic":       req.topic or "General Research",
        "paper_count": len(papers),
        "ideas":       ideas[:5],
        "ai_model":    "llama-3.3-70b-versatile (Groq)",
    }


# ── AI Status ─────────────────────────────────────────────────────────────────

@app.get("/api/ai-status")
def ai_status():
    """Check whether the AI engine is active."""
    active = ai_available()
    return {
        "ai_available": active,
        "provider":     "Groq" if active else None,
        "models": {
            "fast":    "llama-3.1-8b-instant",
            "quality": "llama-3.3-70b-versatile",
        } if active else {},
        "features": ["summarize", "gaps", "chat", "review", "ideas"] if active else [],
        "setup":    "Add GROQ_API_KEY to your .env file" if not active else "Ready",
    }


# ── Full Workflow ─────────────────────────────────────────────────────────────

@app.post("/api/workflow")
def full_workflow(req: SearchRequest):
    """
    Complete workflow: Search → Summarize → Knowledge Graph → Gap Discovery
    Returns everything in one call.
    """
    # 1. Search
    papers = find_papers(req.query, max(1, min(req.max_results, 20)))
    if not papers:
        raise HTTPException(404, "No papers found.")
    for p in papers:
        db.save_paper(p)
    db.log_search(req.query, len(papers))

    # 2. Summarize all
    summaries = []
    for p in papers:
        s = summarize_paper(p)
        db.save_summary(p["id"], s)
        summaries.append({"paper_id": p["id"], "title": p["title"], "summary": s})

    # 3. Generate PPTX for top paper
    pptx_result = None
    if papers:
        top     = papers[0]
        top_sum = summaries[0]["summary"] if summaries else summarize_paper(top)
        path    = generate_pptx(top, top_sum)
        pptx_result = {
            "title":     top["title"],
            "file_path": path,
            "download":  f"/api/present/download/{top['id']}",
        }

    # 4. Knowledge graph
    graph = build_knowledge_graph(papers)

    # 5. Gaps
    gaps = discover_gaps(papers)

    return {
        "query":       req.query,
        "papers":      papers,
        "summaries":   summaries,
        "presentation": pptx_result,
        "knowledge_graph": graph,
        "gaps":        gaps,
        "stats": {
            "papers_found":     len(papers),
            "summaries_created": len(summaries),
            "gaps_found":       len(gaps),
            "graph_nodes":      len(graph["nodes"]),
            "graph_edges":      len(graph["edges"]),
        }
    }


# ── Semantic Search ───────────────────────────────────────────────────────────

@app.post("/api/semantic-search")
def semantic_search_endpoint(req: SemanticSearchRequest):
    """Semantic similarity search over your saved paper library."""
    papers = db.get_papers()
    if not papers:
        raise HTTPException(404, "Library is empty. Search for papers first.")
    top_k = max(1, min(req.top_k, len(papers)))
    results = semantic_search(req.query, papers, top_k)
    return {
        "query":          req.query,
        "semantic_mode":  semantic_available(),
        "count":          len(results),
        "papers":         results,
    }


# ── Similar Papers ────────────────────────────────────────────────────────────

@app.get("/api/similar/{paper_id}")
def similar_papers(paper_id: str, top_k: int = Query(default=5, ge=1, le=20)):
    """Find papers semantically similar to a given paper."""
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found.")
    all_papers = db.get_papers()
    similar = find_similar_papers(paper, all_papers, top_k)
    return {
        "paper_id":    paper_id,
        "paper_title": paper["title"],
        "similar":     similar,
        "count":       len(similar),
    }


# ── Trend Analysis ────────────────────────────────────────────────────────────

@app.get("/api/trends")
def trends():
    """Analyze publication trends and topic evolution across your library."""
    papers = db.get_papers()
    if not papers:
        raise HTTPException(404, "Library is empty. Search for papers first.")
    result = analyze_trends(papers)
    return {"paper_count": len(papers), **result}


# ── BibTeX Export ─────────────────────────────────────────────────────────────

@app.get("/api/bibtex")
def bibtex_all():
    """Export all saved papers as BibTeX."""
    papers = db.get_papers()
    if not papers:
        raise HTTPException(404, "Library is empty.")
    entries = "\n\n".join(paper_to_bibtex(p) for p in papers)
    return JSONResponse(
        content={"bibtex": entries, "count": len(papers)},
        headers={"Content-Disposition": "inline; filename=\"scholargenie.bib\""},
    )


@app.get("/api/bibtex/{paper_id}")
def bibtex_single(paper_id: str):
    """Export a single paper as BibTeX."""
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found.")
    return {"paper_id": paper_id, "bibtex": paper_to_bibtex(paper)}


# ── Library Management ────────────────────────────────────────────────────────

@app.delete("/api/papers/{paper_id}")
def delete_paper(paper_id: str):
    """Remove a single paper from the library."""
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(404, "Paper not found.")
    db.delete_paper(paper_id)
    return {"deleted": True, "paper_id": paper_id, "title": paper["title"]}


@app.delete("/api/papers")
def clear_all_papers():
    """Remove ALL papers and summaries from the library."""
    count = db.stats()["total_papers"]
    db.clear_papers()
    return {"deleted": True, "papers_removed": count}


# ── Stats ─────────────────────────────────────────────────────────────────────

@app.get("/api/stats")
def stats():
    """System statistics."""
    return db.stats()


# ── Agent Pipeline ─────────────────────────────────────────────────────────────

class PipelineRequest(BaseModel):
    query: str

@app.get("/api/debug-pipeline")
def debug_pipeline():
    """Debug: show what agents.py the server is loading."""
    import sys, os, traceback as _tb
    info = {
        "cwd": os.getcwd(),
        "sys_path_0": sys.path[0] if sys.path else "empty",
        "agents_file": None,
        "llm_type": None,
        "error": None,
    }
    try:
        from agents import _build_llm
        import agents as _agents_mod
        info["agents_file"] = getattr(_agents_mod, "__file__", "unknown")
        llm = _build_llm()
        info["llm_type"] = type(llm).__name__
        info["llm_ready"] = llm is not None
    except Exception as e:
        info["error"] = _tb.format_exc()
    return info

@app.post("/api/pipeline")
def agent_pipeline(req: PipelineRequest):
    """
    Run the full ScholarGenie CrewAI agent pipeline.
    Sequence: ResearchAgent → AnalysisAgent → GapFinderAgent → PresentationAgent
    """
    if not req.query.strip():
        raise HTTPException(400, "Query cannot be empty.")
    try:
        from agents import run_pipeline, CREWAI_AVAILABLE
        if not CREWAI_AVAILABLE:
            raise HTTPException(503, "CrewAI not installed. Run: pip install crewai")
        result = run_pipeline(req.query.strip())
        if result.get("status") == "error":
            raise HTTPException(500, result.get("message", "Pipeline failed."))
        return result
    except HTTPException:
        raise
    except Exception as e:
        import traceback as _tb
        print("PIPELINE ERROR TRACEBACK:", flush=True)
        print(_tb.format_exc(), flush=True)
        raise HTTPException(500, f"Pipeline error: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# STARTUP
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import uvicorn
    _get_groq()   # warm up AI client at startup
    ai = "✓ Groq AI Active (Llama3)" if ai_available() else "✗ No AI Key  (add GROQ_API_KEY to .env)"
    print("\n" + "="*60)
    print("  ScholarGenie API Server  v2.0")
    print("="*60)
    print(f"  API:      http://localhost:8000")
    print(f"  Docs:     http://localhost:8000/docs")
    print(f"  Database: {DB_PATH}")
    print(f"  AI:       {ai}")
    print("="*60 + "\n")
    uvicorn.run("api:app", host="0.0.0.0", port=8000, reload=False)
